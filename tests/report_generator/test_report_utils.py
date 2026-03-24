#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from docx import Document
from pptx.oxml.text import CT_TextParagraph
# noinspection PyProtectedMember
from pptx.text.text import _Paragraph

from report_generator.generator import report_utils

class TestReportUtils:

    def test_merge_similar_runs(self):
        p = _Paragraph(CT_TextParagraph(), None)
        r1 = p.add_run()
        r1.text = "aap"
        f1 = r1.font
        f1.bold = True

        r2 = p.add_run()
        r2.text = "noot"
        f2 = r2.font
        f2.bold = True

        report_utils.pptx.merge_runs_with_same_formatting(p)

        assert len(p.runs) == 1
        assert p.text == "aapnoot"

    def test_do_not_merge_different_runs(self):
        p = _Paragraph(CT_TextParagraph(), None)
        r1 = p.add_run()
        r2 = p.add_run()
        r1.text = "aap"
        r2.text = "noot"
        f1 = r1.font
        f1.bold = True
        f2 = r2.font
        f2.bold = False

        report_utils.pptx.merge_runs_with_same_formatting(p)

        assert len(p.runs) == 2
        assert p.runs[0].text == "aap"
        assert p.runs[1].text == "noot"



class TestPptxUpdateParagraphWordBoundaries:
    """Test word boundary behavior in PowerPoint placeholder replacement."""

    def test_substring_not_replaced_in_longer_placeholder(self):
        """Verify FOO does not replace FOOBAR when only FOOBAR is present."""
        p = _Paragraph(CT_TextParagraph(), None)
        r = p.add_run()
        r.text = "The value is PLACEHOLDER_FULL"

        report_utils.pptx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "The value is PLACEHOLDER_FULL"

    def test_exact_match_is_replaced(self):
        """Verify exact placeholder match is replaced correctly."""
        p = _Paragraph(CT_TextParagraph(), None)
        r = p.add_run()
        r.text = "The value is PLACEHOLDER"

        report_utils.pptx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "The value is REPLACED"

    def test_placeholder_with_punctuation_is_replaced(self):
        """Verify placeholder followed by punctuation is replaced."""
        p = _Paragraph(CT_TextParagraph(), None)
        r = p.add_run()
        r.text = "Values: PLACEHOLDER, PLACEHOLDER. PLACEHOLDER!"

        report_utils.pptx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "Values: REPLACED, REPLACED. REPLACED!"

    def test_multiple_runs_selects_correct_run_with_word_boundary(self):
        """Verify correct run is selected when multiple runs contain similar placeholders.
        
        Note: pptx merges runs with same formatting first, so both placeholders end up in one run.
        The word boundary regex ensures only the exact match is replaced.
        """
        p = _Paragraph(CT_TextParagraph(), None)
        
        # First run has longer placeholder
        r1 = p.add_run()
        r1.text = "PLACEHOLDER_FULL "
        
        # Second run has exact match
        r2 = p.add_run()
        r2.text = "PLACEHOLDER"

        report_utils.pptx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        # After merge and replacement, both are in same run
        assert p.text == "PLACEHOLDER_FULL REPLACED"

    def test_placeholder_in_quotes(self):
        """Verify placeholder in quotes is replaced."""
        p = _Paragraph(CT_TextParagraph(), None)
        r = p.add_run()
        r.text = '"PLACEHOLDER"'

        report_utils.pptx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == '"REPLACED"'

    def test_placeholder_in_parentheses(self):
        """Verify placeholder in parentheses is replaced."""
        p = _Paragraph(CT_TextParagraph(), None)
        r = p.add_run()
        r.text = "(PLACEHOLDER)"

        report_utils.pptx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "(REPLACED)"

    def test_both_placeholders_in_same_run_short_replaced_only(self):
        """Verify when both SHORT and SHORT_LONG are in same run, only SHORT is replaced."""
        p = _Paragraph(CT_TextParagraph(), None)
        r = p.add_run()
        r.text = "Short: PLACEHOLDER, Full: PLACEHOLDER_FULL"

        report_utils.pptx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "Short: REPLACED, Full: PLACEHOLDER_FULL"


class TestDocxUpdateParagraphWordBoundaries:
    """Test word boundary behavior in Word placeholder replacement."""

    def test_substring_not_replaced_in_longer_placeholder(self):
        """Verify FOO does not replace FOOBAR when only FOOBAR is present."""
        doc = Document()
        p = doc.add_paragraph()
        p.add_run("The value is PLACEHOLDER_FULL")

        report_utils.docx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "The value is PLACEHOLDER_FULL"

    def test_exact_match_is_replaced(self):
        """Verify exact placeholder match is replaced correctly."""
        doc = Document()
        p = doc.add_paragraph()
        p.add_run("The value is PLACEHOLDER")

        report_utils.docx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "The value is REPLACED"

    def test_placeholder_with_punctuation_is_replaced(self):
        """Verify placeholder followed by punctuation is replaced."""
        doc = Document()
        p = doc.add_paragraph()
        p.add_run("Values: PLACEHOLDER, PLACEHOLDER. PLACEHOLDER!")

        report_utils.docx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "Values: REPLACED, REPLACED. REPLACED!"

    def test_multiple_runs_selects_correct_run_with_word_boundary(self):
        """Verify correct run is selected when multiple runs contain similar placeholders."""
        doc = Document()
        p = doc.add_paragraph()
        
        # First run has longer placeholder (should not be selected)
        p.add_run("PLACEHOLDER_FULL ")
        
        # Second run has exact match (should be selected)
        p.add_run("PLACEHOLDER")

        report_utils.docx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.runs[0].text == "PLACEHOLDER_FULL "
        assert p.runs[1].text == "REPLACED"

    def test_placeholder_in_quotes(self):
        """Verify placeholder in quotes is replaced."""
        doc = Document()
        p = doc.add_paragraph()
        p.add_run('"PLACEHOLDER"')

        report_utils.docx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == '"REPLACED"'

    def test_placeholder_in_parentheses(self):
        """Verify placeholder in parentheses is replaced."""
        doc = Document()
        p = doc.add_paragraph()
        p.add_run("(PLACEHOLDER)")

        report_utils.docx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "(REPLACED)"

    def test_both_placeholders_in_same_run_short_replaced_only(self):
        """Verify when both SHORT and SHORT_LONG are in same run, only SHORT is replaced."""
        doc = Document()
        p = doc.add_paragraph()
        p.add_run("Short: PLACEHOLDER, Full: PLACEHOLDER_FULL")

        report_utils.docx.update_paragraph(p, "PLACEHOLDER", "REPLACED")

        assert p.text == "Short: REPLACED, Full: PLACEHOLDER_FULL"
