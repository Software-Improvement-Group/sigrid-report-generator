#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import io

from pptx import Presentation
from pptx.util import Inches

from report_generator.generator.placeholders.rendering import pptx as render
from report_generator.generator.placeholders.rendering.pptx import index as pptx_index
from report_generator.generator.utils.constants.sentiment import Sentiment


def _presentation_with_paragraphs(*lines):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        text_frame.add_paragraph().text = line
    return presentation


def _add_slide_with_text(presentation, text):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    textbox.text_frame.paragraphs[0].text = text
    return slide


def _slide_texts(presentation):
    return [
        shape.text_frame.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]


def test_delete_slides_with_placeholder_removes_only_matching_slide():
    presentation = Presentation()
    _add_slide_with_text(presentation, "keep this SYSTEM_PY slide")
    _add_slide_with_text(presentation, "drop this SECURITY_RATING slide")
    _add_slide_with_text(presentation, "keep this one too")

    render.delete_slides_with_placeholder(presentation, "SECURITY_RATING")

    remaining = _slide_texts(presentation)
    assert len(list(presentation.slides)) == 2
    assert "keep this SYSTEM_PY slide" in remaining
    assert "keep this one too" in remaining
    assert all("SECURITY_RATING" not in text for text in remaining)


def test_delete_slides_with_placeholder_is_idempotent():
    presentation = Presentation()
    _add_slide_with_text(presentation, "drop this SECURITY_RATING slide")
    _add_slide_with_text(presentation, "keep this one")

    render.delete_slides_with_placeholder(presentation, "SECURITY_RATING")
    render.delete_slides_with_placeholder(presentation, "SECURITY_RATING")

    assert len(list(presentation.slides)) == 1


def test_delete_slides_with_placeholder_matches_chart_or_table_shape_name():
    presentation = Presentation()
    slide = _add_slide_with_text(presentation, "no key in the text here")
    slide.shapes[0].name = "SECURITY_CHART"
    _add_slide_with_text(presentation, "keep this one")

    render.delete_slides_with_placeholder(presentation, "SECURITY_CHART")

    assert len(list(presentation.slides)) == 1
    assert "keep this one" in _slide_texts(presentation)


def test_delete_slides_with_placeholder_matches_shape_name_nested_in_group():
    presentation = Presentation()
    slide = _add_slide_with_text(presentation, "no key in the text here")
    group = slide.shapes.add_group_shape([slide.shapes[0]])
    group.shapes[0].name = "SECURITY_CHART"
    _add_slide_with_text(presentation, "keep this one")

    render.delete_slides_with_placeholder(presentation, "SECURITY_CHART")

    assert len(list(presentation.slides)) == 1
    assert "keep this one" in _slide_texts(presentation)


def test_delete_slides_with_placeholder_no_match_keeps_all_slides():
    presentation = Presentation()
    _add_slide_with_text(presentation, "slide one")
    _add_slide_with_text(presentation, "slide two")

    render.delete_slides_with_placeholder(presentation, "MISSING_KEY")

    assert len(list(presentation.slides)) == 2


def test_finds_placeholder_in_every_paragraph_of_same_text_frame():
    presentation = _presentation_with_paragraphs(
        "rebuild value of SYSTEM_PY person years",
        "no placeholder here",
        "or SYSTEM_PY person years (PY)",
    )

    paragraphs = render.find_text_in_presentation(presentation, "SYSTEM_PY")

    assert len(paragraphs) == 2
    assert all("SYSTEM_PY" in p.text for p in paragraphs)


def test_update_replaces_placeholder_across_paragraphs():
    presentation = _presentation_with_paragraphs(
        "rebuild value of SYSTEM_PY person years",
        "or SYSTEM_PY person years (PY)",
    )

    paragraphs = render.find_text_in_presentation(presentation, "SYSTEM_PY")
    render.update_many_paragraphs(paragraphs, "SYSTEM_PY", "42.0")

    texts = [p.text for p in render.find_text_in_presentation(presentation, "42.0")]
    assert len(texts) == 2
    assert all("SYSTEM_PY" not in text for text in texts)


def test_find_shapes_with_text_deduplicates_shared_shape():
    presentation = _presentation_with_paragraphs(
        "SYSTEM_PY one",
        "SYSTEM_PY two",
    )

    shapes = render.find_shapes_with_text(presentation, "SYSTEM_PY")

    assert len(shapes) == 1


def test_search_text_is_matched_literally_not_as_regex():
    # "42.0" must not match "4250": the dot is a literal, not a wildcard.
    presentation = _presentation_with_paragraphs("value is 4250 today")

    assert render.find_text_in_presentation(presentation, "42.0") == []


def test_sentiment_color_maps_sentiment_to_color():
    assert render.sentiment_color(Sentiment.POSITIVE) == render.FIVE_STAR_COLOR  # green
    assert render.sentiment_color(Sentiment.NEGATIVE) == render.ONE_STAR_COLOR  # red
    assert render.sentiment_color(Sentiment.NEUTRAL) == render.SIG_BLUE_COLOR  # blue


def _presentation_with_group(text):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text_frame.paragraphs[0].text = text
    group = slide.shapes.add_group_shape([textbox])
    return presentation, group, next(iter(group.shapes))


def _reopened(presentation):
    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return Presentation(buffer)


def _all_texts(presentation):
    return [
        record.text for record in pptx_index.for_presentation(presentation).paragraphs
    ]


def test_paragraph_in_group_keeps_the_ancestor_chain_placeholders_walk_up():
    # Several placeholders locate the shape to move by walking back up from the paragraph:
    # the marker placeholders use four hops, the summary ones two. Those links only exist on
    # proxies built by descending the shape tree, so this pins the traversal to the proxy API.
    presentation, group, inner_shape = _presentation_with_group("nested SYSTEM_PY here")

    for lookup in range(2):  # cold index, then warm index
        paragraphs = render.find_text_in_presentation(presentation, "SYSTEM_PY")

        assert len(paragraphs) == 1, f"lookup {lookup}"
        paragraph = paragraphs[0]
        assert paragraph._parent._parent._parent._parent == group
        assert paragraph._parent._parent == inner_shape


def test_repeated_lookups_return_equal_paragraphs():
    presentation = _presentation_with_paragraphs(
        "rebuild value of SYSTEM_PY person years"
    )

    first = render.find_text_in_presentation(presentation, "SYSTEM_PY")
    second = render.find_text_in_presentation(presentation, "SYSTEM_PY")

    # Proxies are views over the document, so a fresh one compares equal but is not identical.
    assert first == second


def test_writes_through_a_warm_index_reach_the_saved_file():
    presentation = _presentation_with_paragraphs(
        "rebuild value of SYSTEM_PY person years"
    )
    slide = next(iter(presentation.slides))
    table_shape = slide.shapes.add_table(
        1, 1, Inches(1), Inches(3), Inches(4), Inches(1)
    )
    table_shape.table.cell(0, 0).text_frame.paragraphs[0].text = "cell SYSTEM_TC value"
    grouped_textbox = slide.shapes.add_textbox(
        Inches(1), Inches(5), Inches(3), Inches(1)
    )
    grouped_textbox.text_frame.paragraphs[0].text = "grouped SYSTEM_GP value"
    slide.shapes.add_group_shape([grouped_textbox])

    # Warm the index before writing, so the writes go through cached proxies.
    render.find_text_in_presentation(presentation, "SYSTEM_PY")
    for key, value in (
        ("SYSTEM_PY", "11.0"),
        ("SYSTEM_TC", "22.0"),
        ("SYSTEM_GP", "33.0"),
    ):
        render.update_many_paragraphs(
            render.find_text_in_presentation(presentation, key), key, value
        )

    texts = " ".join(_all_texts(_reopened(presentation)))
    assert "11.0" in texts and "22.0" in texts and "33.0" in texts
    assert (
        "SYSTEM_PY" not in texts
        and "SYSTEM_TC" not in texts
        and "SYSTEM_GP" not in texts
    )


def test_remove_shape_drops_the_shape_from_later_lookups():
    presentation = _presentation_with_paragraphs(
        "rebuild value of SYSTEM_PY person years"
    )
    shape = render.find_shapes(presentation, "SYSTEM_PY")[0]

    render.remove_shape(shape)

    assert render.find_text_in_presentation(presentation, "SYSTEM_PY") == []
    assert "SYSTEM_PY" not in " ".join(_all_texts(_reopened(presentation)))


def test_delete_slides_with_placeholder_does_not_rebuild_when_nothing_matches(
    monkeypatch,
):
    # The failure path runs for every failing placeholder, and an unlicensed report fails
    # hundreds without matching a slide, so a no-match call must not discard the index.
    presentation = Presentation()
    _add_slide_with_text(presentation, "keep this SYSTEM_PY slide")
    builds = []
    original = pptx_index.walk.presentation_index
    monkeypatch.setattr(
        pptx_index.walk,
        "presentation_index",
        lambda pres: (builds.append(1), original(pres))[1],
    )
    render.find_text_in_presentation(presentation, "SYSTEM_PY")

    render.delete_slides_with_placeholder(presentation, "ABSENT_KEY")
    render.find_text_in_presentation(presentation, "SYSTEM_PY")

    assert len(builds) == 1


def test_delete_slides_with_placeholder_rebuilds_after_removing_a_slide():
    presentation = Presentation()
    _add_slide_with_text(presentation, "drop this SECURITY_RATING slide")
    _add_slide_with_text(presentation, "keep this SYSTEM_PY slide")
    render.find_text_in_presentation(presentation, "SYSTEM_PY")

    render.delete_slides_with_placeholder(presentation, "SECURITY_RATING")

    assert render.find_text_in_presentation(presentation, "SECURITY_RATING") == []
    assert len(render.find_text_in_presentation(presentation, "SYSTEM_PY")) == 1


def test_find_charts_and_tables_ignore_shapes_nested_in_a_group():
    # Charts and tables are located by shape name among the slide's own shapes only. A nested
    # one is unreachable today; the slide containing it is still dropped on failure.
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        1, 1, Inches(1), Inches(1), Inches(4), Inches(1)
    )
    table_shape.name = "NESTED_TABLE"
    slide.shapes.add_group_shape([table_shape])

    assert render.find_tables(presentation, "NESTED_TABLE") == []
    render.delete_slides_with_placeholder(presentation, "NESTED_TABLE")
    assert len(list(presentation.slides)) == 0


def test_find_shapes_returns_the_slide_level_shape_for_nested_text():
    presentation, group, _inner_shape = _presentation_with_group(
        "nested SYSTEM_PY here"
    )

    shapes = render.find_shapes(presentation, "SYSTEM_PY")

    assert len(shapes) == 1
    assert shapes[0] == group


def test_find_text_in_slide_still_searches_a_slide_removed_from_the_deck():
    presentation = Presentation()
    _add_slide_with_text(presentation, "orphaned SYSTEM_PY slide")
    slide = next(iter(presentation.slides))
    render.find_text_in_presentation(presentation, "SYSTEM_PY")
    id_list = presentation.slides.element
    id_list.remove(next(iter(id_list)))

    # The slide is no longer part of the deck, so it is absent from the index. Walking it is
    # the honest answer; claiming it has no text would hide the caller's mistake.
    assert len(render.find_text_in_slide(slide, "SYSTEM_PY")) == 1
