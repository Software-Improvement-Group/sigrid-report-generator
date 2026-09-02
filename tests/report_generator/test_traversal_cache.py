#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from pptx import Presentation

from report_generator.generator.placeholders.rendering.pptx.index import cache


class _CountingBuilder:
    def __init__(self, value="index"):
        self.calls = 0
        self.value = value

    def __call__(self):
        self.calls += 1
        return self.value


def test_index_is_built_once_then_reused():
    presentation = Presentation()
    builder = _CountingBuilder()

    first = cache.index_for(presentation, builder)
    second = cache.index_for(presentation, builder)

    assert first is second
    assert builder.calls == 1


def test_invalidate_forces_a_rebuild():
    presentation = Presentation()
    builder = _CountingBuilder()

    cache.index_for(presentation, builder)
    cache.invalidate(presentation)
    cache.index_for(presentation, builder)

    assert builder.calls == 2


def test_any_proxy_of_the_same_document_is_an_equivalent_anchor():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    builder = _CountingBuilder()

    cache.index_for(presentation, builder)
    # A slide, a shape or a paragraph must reach the same entry as the presentation itself,
    # because the invalidation and refresh hooks only ever have one of those in hand.
    cache.index_for(slide, builder)

    assert builder.calls == 1


def test_a_different_document_replaces_the_cached_entry():
    first_presentation = Presentation()
    second_presentation = Presentation()
    first_builder = _CountingBuilder("first")
    second_builder = _CountingBuilder("second")

    assert cache.index_for(first_presentation, first_builder) == "first"
    assert cache.index_for(second_presentation, second_builder) == "second"
    # The first document was pushed out by the second, so it has to be rebuilt.
    assert cache.index_for(first_presentation, first_builder) == "first"

    assert first_builder.calls == 2
    assert second_builder.calls == 1


def test_cached_index_returns_none_before_anything_is_built():
    presentation = Presentation()

    assert cache.cached_index(presentation) is None

    cache.index_for(presentation, _CountingBuilder())
    assert cache.cached_index(presentation) == "index"


def test_invalidate_is_safe_for_an_unknown_document():
    cache.invalidate(Presentation())


def test_word_bounded_match_treats_search_text_literally():
    assert cache.matches_word_bounded("value is 42.0 today", "42.0")
    # "." must not act as a wildcard, and a partial word must not match.
    assert not cache.matches_word_bounded("value is 4250 today", "42.0")
    assert not cache.matches_word_bounded("MARKER_MAINT_RATING", "MAINT_RATING")
