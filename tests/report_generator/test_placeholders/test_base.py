#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import logging
import re
from unittest.mock import MagicMock, patch

from docx import Document
from pptx import Presentation

from report_generator.generator.placeholders.implementations import placeholders
from report_generator.generator.placeholders.implementations.base import (
    PlaceholderDocType,
)
from report_generator.generator.report import Report, ReportType
from report_generator.generator.utils.constants import MetricEnum


class TestPlaceholders:
    def test_to_json_name(self):
        class TestMetricEnum(MetricEnum):
            UNIT_SIZE = "UNIT_SIZE"
            DUPLICATION = "DUPLICATION"
            duplication = "duplication"
            UnIt_sIze = "UnIt_sIze"

        assert TestMetricEnum.UNIT_SIZE.to_json_name() == "unitSize"
        assert TestMetricEnum.DUPLICATION.to_json_name() == "duplication"
        assert TestMetricEnum.duplication.to_json_name() == "duplication"
        assert TestMetricEnum.UnIt_sIze.to_json_name() == "unitSize"

    def test_all_placeholder_keys_use_valid_characters(self):
        """All placeholder keys must only contain uppercase letters, digits, and underscores,
        and must contain at least one underscore. Parameterized keys may additionally contain {parameter} tokens."""
        valid_key = re.compile(r"^(?=[A-Z0-9_]*_)[A-Z0-9_]+(\{[A-Z0-9_a-z]+\}[A-Z0-9_]*)*$")
        invalid = [ph.key for ph in placeholders if not valid_key.match(ph.key)]
        assert not invalid, f"Placeholder keys with invalid characters: {invalid}"

    def test_all_placeholders_produce_finds_logs(self, caplog):
        """Verify all placeholders produce 'Finds for' debug logs when resolved."""

        def _get_dummy_value(doc_type):
            """Return appropriate dummy value based on PlaceholderDocType."""
            if doc_type == PlaceholderDocType.TEXT:
                return "Dummy Text"
            elif doc_type == PlaceholderDocType.CHART:
                return {"categories": ["A", "B"], "values": [1, 2]}
            elif doc_type == PlaceholderDocType.TABLE:
                return [["Header"], ["Row1"]]
            elif doc_type == PlaceholderDocType.IMAGE:
                return b"dummy_image_data"
            else:
                return "Dummy Value"

        # Create mock reports for both types
        mock_pptx = MagicMock(spec=Presentation)
        mock_pptx.slides = []
        report_pptx = Report(mock_pptx, ReportType.PRESENTATION)

        mock_docx = MagicMock(spec=Document)
        mock_docx.paragraphs = []
        mock_docx.tables = []
        report_docx = Report(mock_docx, ReportType.DOCUMENT)

        expected_log_entries = set()

        with caplog.at_level(logging.DEBUG):
            for placeholder_cls in placeholders:
                # Determine which report types this placeholder supports
                supports_pptx = placeholder_cls.supports(ReportType.PRESENTATION)
                supports_docx = placeholder_cls.supports(ReportType.DOCUMENT)

                if not supports_pptx and not supports_docx:
                    continue

                # Get dummy value based on doc type
                doc_type = getattr(
                    placeholder_cls, "__doc_type__", PlaceholderDocType.OTHER
                )
                dummy_value = _get_dummy_value(doc_type)

                # Mock the value method
                with patch.object(placeholder_cls, "value", return_value=dummy_value):
                    # Handle parameterized placeholders
                    if placeholder_cls.is_parameterized():
                        # Collect expected keys for all parameters
                        for param in placeholder_cls.allowed_parameters:
                            key = placeholder_cls.key.format(parameter=param)
                            expected_log_entries.add(key)

                        # Resolve once per report type; ParameterizedPlaceholder.resolve
                        # is responsible for iterating over allowed parameters internally.
                        if supports_pptx:
                            try:
                                placeholder_cls.resolve(report_pptx)
                            except Exception:
                                pass  # Some placeholders may fail with mocked data

                        if supports_docx:
                            try:
                                placeholder_cls.resolve(report_docx)
                            except Exception:
                                pass
                    else:
                        # Regular placeholder
                        key = placeholder_cls.key
                        expected_log_entries.add(key)

                        # Resolve for supported report types
                        if supports_pptx:
                            try:
                                placeholder_cls.resolve(report_pptx)
                            except Exception:
                                pass

                        if supports_docx:
                            try:
                                placeholder_cls.resolve(report_docx)
                            except Exception:
                                pass

        # Verify all expected keys appear in debug logs
        log_text = caplog.text
        missing_logs = [
            key for key in expected_log_entries if f"Finds for {key}:" not in log_text
        ]

        assert not missing_logs, (
            f"Missing 'Finds for' logs for {len(missing_logs)} placeholders: {missing_logs[:10]}"
        )
