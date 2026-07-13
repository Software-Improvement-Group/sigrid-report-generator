#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import pytest
from pptx import Presentation
from pptx.util import Inches

from report_generator.generator.domain.portfolio.maintainability_portfolio.statistics import (
    maintainability_portfolio_stats,
)
from report_generator.generator.placeholders.implementations.text.maintainability_portfolio import (
    portfolio_maint_average_delta,
)
from report_generator.generator.placeholders.rendering import pptx as render


@pytest.fixture
def prime_averages():
    """Prime the portfolio stats singleton's caches and clean them up afterwards."""

    def _prime(start_average, end_average):
        maintainability_portfolio_stats.__dict__["statistics"] = {
            "maintainability": {
                "start-average": start_average,
                "end-average": end_average,
            }
        }
        maintainability_portfolio_stats.__dict__.pop("average_delta", None)

    yield _prime

    maintainability_portfolio_stats.__dict__.pop("statistics", None)
    maintainability_portfolio_stats.__dict__.pop("average_delta", None)


def _presentation_with_text(text):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    textbox.text_frame.paragraphs[0].text = text
    return presentation


def _resolve(presentation):
    placeholder = portfolio_maint_average_delta
    placeholder.resolve_pptx(presentation, placeholder.key, placeholder.value)


def _only_run(presentation):
    paragraph = presentation.slides[0].shapes[0].text_frame.paragraphs[0]
    return paragraph.runs[0]


def test_key_is_derived_from_function_name():
    assert portfolio_maint_average_delta.key == "PORTFOLIO_MAINT_AVERAGE_DELTA"


def test_increase_renders_signed_value_in_green(prime_averages):
    prime_averages(3.2, 3.5)
    presentation = _presentation_with_text(portfolio_maint_average_delta.key)

    _resolve(presentation)

    run = _only_run(presentation)
    assert run.text == "+0.30"
    assert run.font.color.rgb == render.FIVE_STAR_COLOR


def test_decrease_renders_signed_value_in_red(prime_averages):
    prime_averages(4.0, 3.99)
    presentation = _presentation_with_text(portfolio_maint_average_delta.key)

    _resolve(presentation)

    run = _only_run(presentation)
    assert run.text == "-0.01"
    assert run.font.color.rgb == render.ONE_STAR_COLOR


def test_unchanged_renders_equals_in_blue(prime_averages):
    prime_averages(3.5, 3.5)
    presentation = _presentation_with_text(portfolio_maint_average_delta.key)

    _resolve(presentation)

    run = _only_run(presentation)
    assert run.text == "="
    assert run.font.color.rgb == render.SIG_BLUE_COLOR
