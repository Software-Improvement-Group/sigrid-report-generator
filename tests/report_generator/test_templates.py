#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import re
import tempfile
from pathlib import Path

from docx import Document
from importlib_resources import files
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class TestTemplates:

    def test_templates_should_not_use_sigsterren_font(self):
        """
        Test that no templates use the Sigsterren font at any point. 

        Placeholders should be written in the Calibri font. For stars, the unicode stars (★★★☆☆) should be used with the Calibri font.
        """
        templates_dir = files("report_generator.presets.templates")
        
        template_files = []
        original_names = {}
        
        for resource in templates_dir.iterdir():
            if resource.is_file() and resource.name.endswith(('.pptx', '.docx')):
                if not resource.name.startswith('~$') and not resource.name.startswith('.'):
                    with tempfile.NamedTemporaryFile(delete=False, suffix=resource.name) as tmp:
                        tmp.write(resource.read_bytes())
                        temp_path = Path(tmp.name)
                        template_files.append(temp_path)
                        original_names[temp_path] = resource.name
        
        assert len(template_files) > 0, "No template files found"
        
        fonts_found = {}
        
        try:
            for template_file in template_files:
                original_name = original_names[template_file]
                
                if template_file.suffix == '.pptx':
                    fonts = self._check_pptx_for_sigsterren_font(template_file)
                    if fonts:
                        fonts_found[original_name] = fonts
                        
                elif template_file.suffix == '.docx':
                    fonts = self._check_docx_for_sigsterren_font(template_file)
                    if fonts:
                        fonts_found[original_name] = fonts
        finally:
            for template_file in template_files:
                try:
                    template_file.unlink()
                except Exception:
                    pass
        
        if fonts_found:
            error_msg = "The following templates use the Sigsterren font:\n\n"
            for template, locations in fonts_found.items():
                error_msg += f"  {template}:\n"
                for location in locations:
                    error_msg += f"    - {location}\n"
            assert False, error_msg
    
    def _check_pptx_for_sigsterren_font(self, file_path: Path) -> list[str]:
        fonts_found = []
        prs = Presentation(str(file_path))
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape_idx, shape in enumerate(slide.shapes, start=1):
                self._check_shape_for_sigsterren(shape, slide_idx, shape_idx, fonts_found)
        
        for master_idx, slide_master in enumerate(prs.slide_masters, start=1):
            for shape in slide_master.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and self._is_sigsterren_font(run.font.name):
                                fonts_found.append(
                                    f"Slide Master {master_idx}: font '{run.font.name}', text: '{run.text}'"
                                )
            
            for layout_idx, slide_layout in enumerate(slide_master.slide_layouts, start=1):
                for shape in slide_layout.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name and self._is_sigsterren_font(run.font.name):
                                    fonts_found.append(
                                        f"Slide Master {master_idx}, Layout {layout_idx}: font '{run.font.name}', text: '{run.text}'"
                                    )
        
        return fonts_found
    
    def _check_shape_for_sigsterren(self, shape, slide_idx: int, shape_idx: int, fonts_found: list[str]) -> None:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and self._is_sigsterren_font(run.font.name):
                        fonts_found.append(
                            f"Slide {slide_idx}, Shape {shape_idx}: font '{run.font.name}', text: '{run.text}'"
                        )
        
        if hasattr(shape, "has_table") and shape.has_table:
            for row_idx, row in enumerate(shape.table.rows, start=1):
                for cell_idx, cell in enumerate(row.cells, start=1):
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and self._is_sigsterren_font(run.font.name):
                                fonts_found.append(
                                    f"Slide {slide_idx}, Table, Row {row_idx}, Cell {cell_idx}: font '{run.font.name}', text: '{run.text}'"
                                )
        
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for grouped_shape in shape.shapes:
                self._check_shape_for_sigsterren(grouped_shape, slide_idx, shape_idx, fonts_found)
    
    def _check_docx_for_sigsterren_font(self, file_path: Path) -> list[str]:
        fonts_found = []
        doc = Document(str(file_path))
        
        for para_idx, para in enumerate(doc.paragraphs, start=1):
            for run in para.runs:
                if run.font.name and self._is_sigsterren_font(run.font.name):
                    fonts_found.append(
                        f"Paragraph {para_idx}: font '{run.font.name}', text: '{run.text}'"
                    )
        
        for table_idx, table in enumerate(doc.tables, start=1):
            for row_idx, row in enumerate(table.rows, start=1):
                for cell_idx, cell in enumerate(row.cells, start=1):
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.font.name and self._is_sigsterren_font(run.font.name):
                                fonts_found.append(
                                    f"Table {table_idx}, Row {row_idx}, Cell {cell_idx}: font '{run.font.name}', text: '{run.text}'"
                                )
        
        for style in doc.styles:
            if hasattr(style, 'font') and style.font and style.font.name:
                if self._is_sigsterren_font(style.font.name):
                    fonts_found.append(
                        f"Style '{style.name}': font '{style.font.name}'"
                    )
        
        return fonts_found
    
    def _is_sigsterren_font(self, font_name: str) -> bool:
        font_name_lower = font_name.lower()
        return 'sigsterren' in font_name_lower or 'sig sterren' in font_name_lower
