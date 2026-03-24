#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import sys
from difflib import Differ

from pptx import Presentation


def _changed_lines(text_current: str, text_reference: str) -> list[str]:
    label_map = {"- ": "reference: ", "+ ": "current:   "}
    return [
        label_map[line[:2]] + line[2:]
        for line in Differ().compare(
            text_reference.splitlines(), text_current.splitlines()
        )
        if line.startswith("+ ") or line.startswith("- ")
    ]


def compare_pptx(current_file, reference_file):
    current_prs = Presentation(current_file)
    reference_prs = Presentation(reference_file)

    differences = []

    differences += compare_slide_count(current_prs, reference_prs)
    differences += compare_slides(current_prs, reference_prs)

    return not differences, differences


def compare_slide_count(current_prs, reference_prs):
    differences = []
    if len(current_prs.slides) != len(reference_prs.slides):
        differences.append(
            f"Slide count mismatch: current has {len(current_prs.slides)} slides, reference has {len(reference_prs.slides)} slides"
        )
    return differences


def compare_slides(current_prs, reference_prs):
    differences = []
    for i, (current_slide, reference_slide) in enumerate(
        zip(current_prs.slides, reference_prs.slides)
    ):
        differences += compare_shape_count(current_slide, reference_slide, i)
        differences += compare_shapes_text(current_slide, reference_slide, i)
        differences += compare_tables(current_slide, reference_slide, i)
        differences += compare_charts(current_slide, reference_slide, i)
    return differences


def compare_shape_count(current_slide, reference_slide, slide_index):
    differences = []
    if len(current_slide.shapes) != len(reference_slide.shapes):
        differences.append(
            f"Slide {slide_index + 1}: Shape count mismatch: "
            f"current has {len(current_slide.shapes)} shapes, reference has {len(reference_slide.shapes)} shapes"
        )
    return differences


def compare_shapes_text(current_slide, reference_slide, slide_index):
    differences = []
    for j, (current_shape, reference_shape) in enumerate(
        zip(current_slide.shapes, reference_slide.shapes)
    ):
        if current_shape.has_text_frame and reference_shape.has_text_frame:
            if current_shape.text != reference_shape.text:
                differences.append(
                    f"Slide {slide_index + 1}, Shape {j + 1}: Text difference:"
                )
                differences.extend(
                    f"    {line}"
                    for line in _changed_lines(current_shape.text, reference_shape.text)
                )
    return differences


def compare_tables(current_slide, reference_slide, slide_index):
    differences = []
    for j, (current_shape, reference_shape) in enumerate(
        zip(current_slide.shapes, reference_slide.shapes)
    ):
        if current_shape.has_table and reference_shape.has_table:
            current_table, reference_table = current_shape.table, reference_shape.table
            if len(current_table.rows) != len(reference_table.rows) or len(
                current_table.columns
            ) != len(reference_table.columns):
                differences.append(
                    f"Slide {slide_index + 1}, Table {j + 1}: Table size mismatch: "
                    f"current has {len(current_table.rows)} rows and {len(current_table.columns)} columns, "
                    f"reference has {len(reference_table.rows)} rows and {len(reference_table.columns)} columns"
                )
            else:
                differences += compare_table_cells(
                    current_table, reference_table, slide_index, j
                )
    return differences


def compare_table_cells(current_table, reference_table, slide_index, table_index):
    differences = []
    for row_idx, (current_row, reference_row) in enumerate(
        zip(current_table.rows, reference_table.rows)
    ):
        for col_idx, (current_cell, reference_cell) in enumerate(
            zip(current_row.cells, reference_row.cells)
        ):
            if current_cell.text != reference_cell.text:
                differences.append(
                    f"Slide {slide_index + 1}, Table {table_index + 1}, Cell ({row_idx + 1}, {col_idx + 1}): Text difference:"
                )
                differences.extend(
                    f"    <{line}>"
                    for line in _changed_lines(current_cell.text, reference_cell.text)
                )
    return differences


def compare_charts(current_slide, reference_slide, slide_index):
    differences = []
    for j, (current_shape, reference_shape) in enumerate(
        zip(current_slide.shapes, reference_slide.shapes)
    ):
        if current_shape.has_chart and reference_shape.has_chart:
            differences += compare_chart(
                current_shape.chart, reference_shape.chart, slide_index, j
            )
    return differences


def compare_chart(current_chart, reference_chart, slide_index, shape_index):
    return compare_chart_series(
        current_chart, reference_chart, slide_index, shape_index
    ) + compare_chart_categories(
        current_chart, reference_chart, slide_index, shape_index
    )


def compare_chart_series(current_chart, reference_chart, slide_index, shape_index):
    differences = []
    current_series, reference_series = (
        list(current_chart.series),
        list(reference_chart.series),
    )
    if len(current_series) != len(reference_series):
        differences.append(
            f"Slide {slide_index + 1}, Chart {shape_index + 1}: Series count mismatch:"
        )
        differences.append(f"    reference: {len(reference_series)}")
        differences.append(f"    current:   {len(current_series)}")
        return differences
    for k, (current_s, reference_s) in enumerate(zip(current_series, reference_series)):
        differences += compare_series(
            current_s, reference_s, slide_index, shape_index, k
        )
    return differences


def compare_series(
    current_series, reference_series, slide_index, shape_index, series_index
):
    prefix = (
        f"Slide {slide_index + 1}, Chart {shape_index + 1}, Series {series_index + 1}"
    )
    differences = []
    if current_series.name != reference_series.name:
        differences.append(f"{prefix}: Name mismatch:")
        differences.append(f"    reference: {reference_series.name!r}")
        differences.append(f"    current:   {current_series.name!r}")
    current_vals, reference_vals = (
        list(current_series.values),
        list(reference_series.values),
    )
    if current_vals != reference_vals:
        differences.append(f"{prefix}: Values mismatch:")
        differences.append(f"    reference: {reference_vals}")
        differences.append(f"    current:   {current_vals}")
    return differences


def compare_chart_categories(current_chart, reference_chart, slide_index, shape_index):
    try:
        reference_cats = [c.label for c in reference_chart.plots[0].categories]
        current_cats = [c.label for c in current_chart.plots[0].categories]
    except (AttributeError, IndexError, TypeError):
        return []  # chart type does not expose categories
    if current_cats == reference_cats:
        return []
    return [
        f"Slide {slide_index + 1}, Chart {shape_index + 1}: Category labels mismatch:",
        f"    reference: {reference_cats}",
        f"    current:   {current_cats}",
    ]


def main():
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <file1.pptx> <file2.pptx>", file=sys.stderr)
        sys.exit(1)
    current_file, reference_file = sys.argv[1], sys.argv[2]

    are_equal, differences = compare_pptx(current_file, reference_file)

    if not are_equal:
        print("Differences found:")
        for diff in differences:
            print(diff)
    else:
        print("No differences found.")


if __name__ == "__main__":
    main()
