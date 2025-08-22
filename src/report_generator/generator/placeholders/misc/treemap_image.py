#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from typing import Callable

from pptx import Presentation

from report_generator.generator import report_utils
from report_generator.generator.data_models import maintainability_portfolio_data, security_ratings_portfolio_data, architecture_portfolio_data
from report_generator.generator.placeholders.base import Placeholder
import plotly.express as px
import io
import logging
from pptx.util import Inches

class _AbstractTreemapPlaceholder(Placeholder):
    SIG_BLUE = f"#{report_utils.pptx.SIG_BLUE}"
    NA_STAR_COLOR = f"#{report_utils.pptx.NA_STAR_COLOR}"

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        slides = report_utils.pptx.identify_specific_slide(presentation, key)
        if len(slides) == 0:
            return

        for slide in slides:
            shapes = report_utils.pptx.find_shapes_with_text_in_slide(slide, key)
            for shape in shapes:
                data = value_cb()
                cls.create_and_add_treemap_image_to_slide(shape, slide, data)
    
    @staticmethod
    def determine_rating_color(rating):
        return f"#{report_utils.pptx.determine_rating_color(rating)}"

    @staticmethod    
    def test_code_ratio_color(rating):
        return f"#{report_utils.pptx.test_code_ratio_color(rating)}"

    @staticmethod
    def interpolate_color(colors, t):
        return f"#{report_utils.pptx.interpolate_color(colors, t)}"
    
    @staticmethod
    def normalize_clamped(min_val, max_val, val):
        if min_val == max_val:
            return 0
        return max(0, min(1, (val - min_val) / (max_val - min_val)))

    @staticmethod
    def create_and_add_treemap_image_to_slide(shape_placeholder, slide, data):
        pos_left = shape_placeholder.left.inches
        pos_top = shape_placeholder.top.inches
        pos_width = shape_placeholder.width.inches
        pos_height = shape_placeholder.height.inches

        fig = px.treemap(names=data['names'], parents=data['parents'], values=data['values'], color=data['color'], color_discrete_map=data['color_mapping'])
        fig.update_traces(root_color='rgba(250, 250, 250, 1)')
        fig.update_layout(margin = dict(t=0, l=0, r=0, b=0))
        
        img_bytes = fig.to_image(format="jpg", width=pos_width*2*96, height=pos_height*2*96)
        
        slide.shapes.add_picture(io.BytesIO(img_bytes),
            left=Inches(pos_left), top=Inches(pos_top),
            width=Inches(pos_width), height=Inches(pos_height))

        el = shape_placeholder.element
        el.getparent().remove(el)


class _AbstractPortfolioTreemapPlaceholder(_AbstractTreemapPlaceholder):
    @staticmethod
    def create_portfolio():
        res = {}
        system_names = maintainability_portfolio_data.system_names
        for system_name in system_names:
            md = maintainability_portfolio_data.find_system_metadata(system_name)
            if not md['active'] or md['isDevelopmentOnly']:
                continue

            res[system_name] = {
                'metadata' : md,
                'start_date_data' : maintainability_portfolio_data.start_snapshot(system_name),
                'end_date_data' : maintainability_portfolio_data.end_snapshot(system_name)
            }
        return res
    
    @staticmethod
    def create_blank_portfolio_treemap():
        names = []
        parents = []
        tracking = []

        portfolio = _AbstractPortfolioTreemapPlaceholder.create_portfolio()
        for s in portfolio.values():
            m = s['metadata']
            if m['teamNames'][0] not in names:
                names.append(m['teamNames'][0])
                parents.append("")
                tracking.append("")
            if m['displayName'] is not None:
                names.append(m['displayName'])
            else:
                names.append(m['systemName'])
            parents.append(m['teamNames'][0])
            tracking.append(m['systemName'])

        return {
            'names' : names,
            'parents' : parents,
            'tracking' : tracking
        }
    
    @staticmethod
    def create_treemap_values(portfolio, treemap_data):
        values = [None] * len(treemap_data['tracking'])
        for i, t in enumerate(treemap_data['tracking']):
            cur_parent = treemap_data['parents'][i]
            if cur_parent == "":
                values[i] = 0
                if t not in treemap_data['color_mapping'].keys():
                    treemap_data['color_mapping'][t] = _AbstractPortfolioTreemapPlaceholder.SIG_BLUE
                continue
            if portfolio[t]['end_date_data'] is None:
                values[i] = 0
            values[i] = portfolio[t]['end_date_data']['volumeInPersonMonths']
        return values


class MaintainabilityPortfolioTreemapPlaceholder(_AbstractPortfolioTreemapPlaceholder):
    """Creates a portfolio treemap where the color is determined by the maintainability rating of the individual systems."""

    key = "PORTFOLIO_PERIOD_MAINTAINABILITY"

    @classmethod
    def value(cls, parameter=None):
        portfolio = _AbstractPortfolioTreemapPlaceholder.create_portfolio()
        treemap = _AbstractPortfolioTreemapPlaceholder.create_blank_portfolio_treemap()

        treemap['color_mapping'] = {}
        for t in portfolio.keys():
            maintainability_rating = portfolio[t]['end_date_data']['maintainability']
            treemap['color_mapping'][t] = _AbstractTreemapPlaceholder.determine_rating_color(maintainability_rating) if maintainability_rating is not None else _AbstractTreemapPlaceholder.NA_STAR_COLOR

        treemap['values'] = _AbstractPortfolioTreemapPlaceholder.create_treemap_values(portfolio, treemap)
        treemap['color'] = treemap['tracking']
        return treemap

    
class MaintainabilityChangePortfolioTreemapPlaceholder(_AbstractPortfolioTreemapPlaceholder):
    """Creates a portfolio treemap where the color is determined by the change in maintainability rating of the individual systems during the specified period."""

    key = "PORTFOLIO_PERIOD_MAINTAINABILITY_CHANGE"

    @classmethod
    def value(cls, parameter=None):
        portfolio = _AbstractPortfolioTreemapPlaceholder.create_portfolio()
        treemap = _AbstractPortfolioTreemapPlaceholder.create_blank_portfolio_treemap()

        differences = {}
        for entry in portfolio.keys():
            if portfolio[entry]['start_date_data']['maintainabilityDate'] == portfolio[entry]['end_date_data']['maintainabilityDate']:
                differences[entry] = None
            else:
                differences[entry] = portfolio[entry]['end_date_data']['maintainability']-portfolio[entry]['start_date_data']['maintainability']
        diff_min = min([x for x in differences.values() if x is not None])
        diff_max = max([x for x in differences.values() if x is not None])

        treemap['color_mapping'] = {}
        for i, entry in enumerate(portfolio.keys()):
            diff = differences[entry]
            if diff is None:
                treemap['color_mapping'][entry] = _AbstractPortfolioTreemapPlaceholder.NA_STAR_COLOR
            elif diff < 0:
                t = _AbstractPortfolioTreemapPlaceholder.normalize_clamped(0, abs(diff_min), abs(diff))
                treemap['color_mapping'][entry] = _AbstractTreemapPlaceholder.interpolate_color(report_utils.pptx.MAINTAINABILITY_NEG_CHANGE_RANGE, t)
            else:
                t = _AbstractPortfolioTreemapPlaceholder.normalize_clamped(0, diff_max, diff)
                treemap['color_mapping'][entry] = _AbstractTreemapPlaceholder.interpolate_color(report_utils.pptx.MAINTAINABILITY_POS_CHANGE_RANGE, t)

        treemap['values'] = _AbstractPortfolioTreemapPlaceholder.create_treemap_values(portfolio, treemap)
        treemap['color'] = treemap['tracking']
        return treemap
    

class VolumeChangePortfolioTreemapPlaceholder(_AbstractPortfolioTreemapPlaceholder):
    """Creates a portfolio treemap where the color is determined by the change in volume change (effort) of the individual systems during the specified period."""

    key = "PORTFOLIO_PERIOD_VOLUME_CHANGE"

    @classmethod
    def value(cls, parameter=None):
        portfolio = _AbstractPortfolioTreemapPlaceholder.create_portfolio()
        treemap = _AbstractPortfolioTreemapPlaceholder.create_blank_portfolio_treemap()

        differences = {}
        for entry in portfolio.keys():
            if portfolio[entry]['start_date_data']['maintainabilityDate'] == portfolio[entry]['end_date_data']['maintainabilityDate']:
                differences[entry] = None
            else:
                differences[entry] = portfolio[entry]['end_date_data']['volumeInPersonMonths']-portfolio[entry]['start_date_data']['volumeInPersonMonths']
        diff_min = min([x for x in differences.values() if x is not None])
        diff_max = max([x for x in differences.values() if x is not None])

        treemap['color_mapping'] = {}
        for i, entry in enumerate(portfolio.keys()):
            diff = differences[entry]
            if diff is None:
                treemap['color_mapping'][entry] = _AbstractPortfolioTreemapPlaceholder.NA_STAR_COLOR
            elif diff < 0:
                t = _AbstractPortfolioTreemapPlaceholder.normalize_clamped(0, abs(diff_min), abs(diff))
                treemap['color_mapping'][entry] = _AbstractTreemapPlaceholder.interpolate_color(report_utils.pptx.VOLUME_NEG_CHANGE_RANGE, t)
            else:
                t = _AbstractPortfolioTreemapPlaceholder.normalize_clamped(0, diff_max, diff)
                treemap['color_mapping'][entry] = _AbstractTreemapPlaceholder.interpolate_color(report_utils.pptx.VOLUME_POS_CHANGE_RANGE, t)

        treemap['values'] = _AbstractPortfolioTreemapPlaceholder.create_treemap_values(portfolio, treemap)
        treemap['color'] = treemap['tracking']
        return treemap
    

class TestCodePortfolioTreemapPlaceholder(_AbstractPortfolioTreemapPlaceholder):
    """Creates a portfolio treemap where the color is determined by the test-to-production code ratio of the individual systems."""

    key = "PORTFOLIO_PERIOD_TEST_CODE"

    @classmethod
    def value(cls, parameter=None):
        portfolio = _AbstractPortfolioTreemapPlaceholder.create_portfolio()
        treemap = _AbstractPortfolioTreemapPlaceholder.create_blank_portfolio_treemap()

        treemap['color_mapping'] = {}
        for t in portfolio.keys():
            test_code_ratio = portfolio[t]['end_date_data']['testCodeRatio']
            treemap['color_mapping'][t] = _AbstractTreemapPlaceholder.test_code_ratio_color(float(test_code_ratio)) if test_code_ratio is not None else _AbstractTreemapPlaceholder.NA_STAR_COLOR

        treemap['values'] = _AbstractPortfolioTreemapPlaceholder.create_treemap_values(portfolio, treemap)
        treemap['color'] = treemap['tracking']
        return treemap
    
    
class SecurityRatingsPortfolioTreemapPlaceholder(_AbstractPortfolioTreemapPlaceholder):
    """Creates a portfolio treemap where the color is determined by the security rating of the individual systems."""

    key = "PORTFOLIO_PERIOD_SECURITY_RATINGS"

    @classmethod
    def value(cls, parameter=None):
        portfolio = _AbstractPortfolioTreemapPlaceholder.create_portfolio()
        treemap = _AbstractPortfolioTreemapPlaceholder.create_blank_portfolio_treemap()

        treemap['color_mapping'] = dict.fromkeys(portfolio.keys(), _AbstractTreemapPlaceholder.NA_STAR_COLOR)
        for t in security_ratings_portfolio_data.system_names:
            security_rating = security_ratings_portfolio_data.end_snapshot(t)['rating']
            if security_rating is None:
                logging.debug(f"Cannot find end snapshot for {t}")
                continue
            treemap['color_mapping'][t] = _AbstractTreemapPlaceholder.determine_rating_color(float(security_rating)) if security_rating is not None else _AbstractTreemapPlaceholder.NA_STAR_COLOR

        treemap['values'] = _AbstractPortfolioTreemapPlaceholder.create_treemap_values(portfolio, treemap)
        treemap['color'] = treemap['tracking']
        return treemap


class ArchitecturePortfolioTreemapPlaceholder(_AbstractPortfolioTreemapPlaceholder):
    """Creates a portfolio treemap where the color is determined by the architecture quality rating of the individual systems."""

    key = "PORTFOLIO_PERIOD_ARCHITECTURE"

    @classmethod
    def value(cls, parameter=None):
        portfolio = _AbstractPortfolioTreemapPlaceholder.create_portfolio()
        treemap = _AbstractPortfolioTreemapPlaceholder.create_blank_portfolio_treemap()

        treemap['color_mapping'] = dict.fromkeys(portfolio.keys(), _AbstractTreemapPlaceholder.NA_STAR_COLOR)
        for t in architecture_portfolio_data.system_names:
            architecture_rating = architecture_portfolio_data.end_snapshot(t)['ratings']['architecture']
            if architecture_rating is None:
                logging.debug(f"Cannot find end snapshot for {t}")
                continue
            treemap['color_mapping'][t] = _AbstractTreemapPlaceholder.determine_rating_color(float(architecture_rating)) if architecture_rating is not None else _AbstractTreemapPlaceholder.NA_STAR_COLOR

        treemap['values'] = _AbstractPortfolioTreemapPlaceholder.create_treemap_values(portfolio, treemap)
        treemap['color'] = treemap['tracking']
        return treemap