#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from typing import Callable

from pptx import Presentation

from report_generator.generator import report_utils
from report_generator.generator.data_models import security_dashboard_findings_portfolio_data
from report_generator.generator.data_models import security_dashboard_resolution_times_portfolio_data
from report_generator.generator.data_models import maintainability_portfolio_data
from report_generator.generator.placeholders.base import Placeholder
import plotly.graph_objects as go
import io
import logging
from pptx.util import Inches
from datetime import datetime

class _AbstractChartImagePlaceholder(Placeholder):
    SIG_BLUE_COLOR = f"#{report_utils.pptx.SIG_BLUE_COLOR}"
    NA_STAR_COLOR = f"#{report_utils.pptx.NA_STAR_COLOR}"

    DASHBOARD_EXISTING_FINDINGS_COLOR = f"#{report_utils.pptx.DASHBOARD_EXISTING_FINDINGS_COLOR}"
    DASHBOARD_NEW_FINDINGS_COLOR = f"#{report_utils.pptx.DASHBOARD_NEW_FINDINGS_COLOR}"
    DASHBOARD_RESOLVED_FINDINGS_COLOR = f"#{report_utils.pptx.DASHBOARD_RESOLVED_FINDINGS_COLOR}"

    DASHBOARD_RESOLUTION_NO_RISK_COLOR = f"#{report_utils.pptx.DASHBOARD_RESOLUTION_NO_RISK_COLOR}"
    DASHBOARD_RESOLUTION_LOW_RISK_COLOR = f"#{report_utils.pptx.DASHBOARD_RESOLUTION_LOW_RISK_COLOR}"
    DASHBOARD_RESOLUTION_MEDIUM_RISK_COLOR = f"#{report_utils.pptx.DASHBOARD_RESOLUTION_MEDIUM_RISK_COLOR}"
    DASHBOARD_RESOLUTION_HIGH_RISK_COLOR = f"#{report_utils.pptx.DASHBOARD_RESOLUTION_HIGH_RISK_COLOR}"

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        slides = report_utils.pptx.identify_specific_slide(presentation, key)
        if len(slides) == 0:
            return

        for slide in slides:
            shapes = report_utils.pptx.find_shapes_with_text_in_slide(slide, key)
            for shape in shapes:
                fig = value_cb()
                cls.create_and_add_treemap_image_to_slide(shape, slide, fig)
    
    @staticmethod
    def create_and_add_treemap_image_to_slide(shape_placeholder, slide, fig):
        pos_left = shape_placeholder.left.inches
        pos_top = shape_placeholder.top.inches
        pos_width = shape_placeholder.width.inches
        pos_height = shape_placeholder.height.inches

        fig.update_layout(
            margin = dict(t=0, l=0, r=0, b=0),
            plot_bgcolor="rgba(0,0,0,0)",
            paper_bgcolor="rgba(0,0,0,0)"
        )
        
        img_bytes = fig.to_image(format="png", width=pos_width*2*96, height=pos_height*2*96)
        
        slide.shapes.add_picture(io.BytesIO(img_bytes),
            left=Inches(pos_left), top=Inches(pos_top),
            width=Inches(pos_width), height=Inches(pos_height))

        el = shape_placeholder.element
        el.getparent().remove(el)


class _AbstractSecurityDashboardPlaceholder(_AbstractChartImagePlaceholder):
    LAYOUT = go.Layout(
        xaxis={
            'showline' : True,
            'linewidth' : 2,
            'linecolor' : '#6E7078',
            'type': 'category',
            'categoryorder': 'array',
            'tickmode' : 'array'
        },
        yaxis={
            'showgrid' : True,
            'gridwidth' : 2,
            'gridcolor' : '#E0E4EF'
        },
        legend={
            'orientation' : 'h',
            'yanchor' : 'top',
            'xanchor' : 'center',
            'y' : -0.05,
            'x' : 0.5,
            'traceorder' : 'normal'
        },
        barmode='stack'
    )

    @staticmethod
    def transform_date_labels_to_months(dates):
        return [datetime.strptime(x, "%Y-%m-%d").strftime("%b") for x in dates]

class _AbstractSecurityDashboardFindingsPlaceholder(_AbstractSecurityDashboardPlaceholder):
    @staticmethod
    def create_portfolio():
        res = {"CRITICAL" : {}, "HIGH" : {}, "MEDIUM" : {}, "LOW" : {}}
        for system in security_dashboard_findings_portfolio_data.data['systems']:
            md = maintainability_portfolio_data.find_system_metadata(system['system'])
            if not md or not md['active'] or md['isDevelopmentOnly']:
                continue
            for ratio in system['findingRatio']:
                month = ratio['month']
                for severity in res.keys():
                    if month not in res[severity].keys():
                        res[severity][month] = {"resolved": 0, "existing": 0, "new": 0}
                    for status in res[severity][month].keys():
                        res[severity][month][status] += ratio['severities'][severity][status]
        return res
    
    @staticmethod
    def create_dashboard_with_severity(severity):
        portfolio = _AbstractSecurityDashboardFindingsPlaceholder.create_portfolio()[severity]

        y_values_new = [portfolio[k]['new'] for k in portfolio.keys()]
        y_values_existing = [portfolio[k]['existing'] for k in portfolio.keys()]
        y_values_resolved = [portfolio[k]['resolved'] for k in portfolio.keys()]
        open_findings_text_values = [x + y for x, y in zip(y_values_new, y_values_existing)]
        data = [
            go.Bar(
                x=list(portfolio.keys()),
                y=y_values_new,
                name="New",
                marker_color=_AbstractChartImagePlaceholder.DASHBOARD_NEW_FINDINGS_COLOR,
                offsetgroup="open"
            ),
            go.Bar(
                x=list(portfolio.keys()),
                y=y_values_existing,
                name="Existing",
                marker_color=_AbstractChartImagePlaceholder.DASHBOARD_EXISTING_FINDINGS_COLOR,
                offsetgroup="open",
                textposition="outside",
                text=open_findings_text_values
            ),
            go.Bar(
                x=list(portfolio.keys()),
                y=y_values_resolved,
                name="Resolved",
                marker_color=_AbstractChartImagePlaceholder.DASHBOARD_RESOLVED_FINDINGS_COLOR,
                offsetgroup="closed",
                textposition="outside",
                text=y_values_resolved
            ),
        ]

        layout = _AbstractSecurityDashboardPlaceholder.LAYOUT
        layout.xaxis.update({
            'categoryarray' : list(portfolio.keys()),
            'tickvals' : list(portfolio.keys()),
            'ticktext' : _AbstractSecurityDashboardPlaceholder.transform_date_labels_to_months(portfolio.keys())
        })

        return go.Figure(data=data, layout=layout)
    

class _AbstractSecurityDashboardResolutionTimesPlaceholder(_AbstractSecurityDashboardPlaceholder):
    LEGEND_ENTRIES_PER_SEVERITY = {
        "CRITICAL" : {'noRisk' : "at most 7 days", "lowRisk" : "between 7-14 days", "mediumRisk" : "between 14-30 days", "highRisk" : "at least 30 days"},
        "HIGH" : {'noRisk' : "at most 14 days", "lowRisk" : "between 14-30 days", "mediumRisk" : "between 30-180 days", "highRisk" : "at least 180 days"},
        "MEDIUM" : {'noRisk' : "at most 30 days", "lowRisk" : "between 30-180 days", "mediumRisk" : "between 180-365 days", "highRisk" : "at least 365 days"},
        "LOW" : {'noRisk' : "at most 180 days", "lowRisk" : "between 180-365 days", "mediumRisk" : "between 1-2 years", "highRisk" : "at least 2 years"}
    }
    @staticmethod
    def create_portfolio():
        res = {"CRITICAL" : {}, "HIGH" : {}, "MEDIUM" : {}, "LOW" : {}}
        for system in security_dashboard_resolution_times_portfolio_data.data['systems']:
            md = maintainability_portfolio_data.find_system_metadata(system['system'])
            if not md or not md['active'] or md['isDevelopmentOnly']:
                continue
            for ratio in system['resolutionTimes']:
                month = ratio['month']
                for severity in res.keys():
                    if month not in res[severity].keys():
                        res[severity][month] = {"noRisk": 0, "lowRisk": 0, "mediumRisk": 0, "highRisk" : 0}
                    for status in res[severity][month].keys():
                        res[severity][month][status] += ratio['severities'][severity][status]
        return res
    
    @staticmethod
    def create_dashboard_with_severity(severity):
        portfolio = _AbstractSecurityDashboardResolutionTimesPlaceholder.create_portfolio()[severity]
        legend_entries = _AbstractSecurityDashboardResolutionTimesPlaceholder.LEGEND_ENTRIES_PER_SEVERITY[severity]

        y_values_no_risk = [portfolio[k]['noRisk'] for k in portfolio.keys()]
        y_values_low_risk = [portfolio[k]['lowRisk'] for k in portfolio.keys()]
        y_values_medium_risk = [portfolio[k]['mediumRisk'] for k in portfolio.keys()]
        y_values_high_risk = [portfolio[k]['highRisk'] for k in portfolio.keys()]
        text_values = [x + y for x, y in zip(y_values_no_risk, y_values_low_risk)]
        text_values = [x + y for x, y in zip(text_values, y_values_medium_risk)]
        text_values = [x + y for x, y in zip(text_values, y_values_high_risk)]
        data = [
            go.Bar(
                x=list(portfolio.keys()),
                y=y_values_no_risk,
                name=legend_entries['noRisk'],
                marker_color=_AbstractChartImagePlaceholder.DASHBOARD_RESOLUTION_NO_RISK_COLOR
            ),
            go.Bar(
                x=list(portfolio.keys()),
                y=y_values_low_risk,
                name=legend_entries['lowRisk'],
                marker_color=_AbstractChartImagePlaceholder.DASHBOARD_RESOLUTION_LOW_RISK_COLOR
            ),
            go.Bar(
                x=list(portfolio.keys()),
                y=y_values_medium_risk,
                name=legend_entries['mediumRisk'],
                marker_color=_AbstractChartImagePlaceholder.DASHBOARD_RESOLUTION_MEDIUM_RISK_COLOR
            ),
            go.Bar(
                x=list(portfolio.keys()),
                y=y_values_high_risk,
                name=legend_entries['highRisk'],
                marker_color=_AbstractChartImagePlaceholder.DASHBOARD_RESOLUTION_HIGH_RISK_COLOR,
                textposition="outside",
                text=text_values
            ),
        ]

        layout = _AbstractSecurityDashboardPlaceholder.LAYOUT
        layout.xaxis.update({
            'categoryarray' : list(portfolio.keys()),
            'tickvals' : list(portfolio.keys()),
            'ticktext' : _AbstractSecurityDashboardPlaceholder.transform_date_labels_to_months(portfolio.keys())
        })

        return go.Figure(data=data, layout=layout)


class SecurityDashboardCriticalFindingsPlaceholder(_AbstractSecurityDashboardFindingsPlaceholder):
    """Creates a portfolio bar chart depicting the number of new, existing, and resolved critical security findings of the last 12 months (counting back from the <end_date>)"""

    key = "PORTFOLIO_PERIOD_SECURITY_DASHBOARD_CRITICAL_FINDINGS"

    @classmethod
    def value(cls, parameter=None):
        return _AbstractSecurityDashboardFindingsPlaceholder.create_dashboard_with_severity("CRITICAL")


class SecurityDashboardHighFindingsPlaceholder(_AbstractSecurityDashboardFindingsPlaceholder):
    """Creates a portfolio bar chart depicting the number of new, existing, and resolved high security findings of the last 12 months (counting back from the <end_date>)"""

    key = "PORTFOLIO_PERIOD_SECURITY_DASHBOARD_HIGH_FINDINGS"

    @classmethod
    def value(cls, parameter=None):
        return _AbstractSecurityDashboardFindingsPlaceholder.create_dashboard_with_severity("HIGH")


class SecurityDashboardMediumFindingsPlaceholder(_AbstractSecurityDashboardFindingsPlaceholder):
    """Creates a portfolio bar chart depicting the number of new, existing, and resolved medium security findings of the last 12 months (counting back from the <end_date>)"""

    key = "PORTFOLIO_PERIOD_SECURITY_DASHBOARD_MEDIUM_FINDINGS"

    @classmethod
    def value(cls, parameter=None):
        return _AbstractSecurityDashboardFindingsPlaceholder.create_dashboard_with_severity("MEDIUM")
    

class SecurityDashboardCriticalResolutionTimesPlaceholder(_AbstractSecurityDashboardResolutionTimesPlaceholder):
    """Creates a portfolio bar chart depicting the resolution times of critical security findings of the last 12 months (counting back from the <end_date>)"""

    key = "PORTFOLIO_PERIOD_SECURITY_DASHBOARD_CRITICAL_RESOLUTION_TIMES"

    @classmethod
    def value(cls, parameter=None):
        return _AbstractSecurityDashboardResolutionTimesPlaceholder.create_dashboard_with_severity("CRITICAL")
    

class SecurityDashboardHighResolutionTimesPlaceholder(_AbstractSecurityDashboardResolutionTimesPlaceholder):
    """Creates a portfolio bar chart depicting the resolution times of high security findings of the last 12 months (counting back from the <end_date>)"""

    key = "PORTFOLIO_PERIOD_SECURITY_DASHBOARD_HIGH_RESOLUTION_TIMES"

    @classmethod
    def value(cls, parameter=None):
        return _AbstractSecurityDashboardResolutionTimesPlaceholder.create_dashboard_with_severity("HIGH")
    

class SecurityDashboardMediumResolutionTimesPlaceholder(_AbstractSecurityDashboardResolutionTimesPlaceholder):
    """Creates a portfolio bar chart depicting the resolution times of medium security findings of the last 12 months (counting back from the <end_date>)"""

    key = "PORTFOLIO_PERIOD_SECURITY_DASHBOARD_MEDIUM_RESOLUTION_TIMES"

    @classmethod
    def value(cls, parameter=None):
        return _AbstractSecurityDashboardResolutionTimesPlaceholder.create_dashboard_with_severity("MEDIUM")