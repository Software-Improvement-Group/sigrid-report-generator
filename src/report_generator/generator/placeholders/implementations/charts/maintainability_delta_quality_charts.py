#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from collections.abc import Callable

from pptx.chart.data import CategoryChartData
from pptx.presentation import Presentation

from report_generator.generator.domain import (
    DELTA_QUALITY_METRICS,
    maintainability_delta_quality_system_by_type,
)
from report_generator.generator.domain.system.maintainability_delta_quality import (
    DeltaColumn,
)
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.formatting.formatters import (
    normalize_percentages,
)
from report_generator.generator.placeholders.implementations.base import (
    MultiParameterList,
    ParameterizedPlaceholder,
    PlaceholderDocType,
)
from report_generator.generator.utils.constants.metrics import DeltaType, MaintMetric

# User-facing chart category label for each comparison column (presentation concern).
_COLUMN_LABELS = {
    DeltaColumn.TOTAL_BEFORE: "Total code (Before)",
    DeltaColumn.NEW_AFTER: "New code (After)",
    DeltaColumn.CHANGED_BEFORE: "Changed code (Before)",
    DeltaColumn.CHANGED_AFTER: "Changed code (After)",
    DeltaColumn.DELETED_BEFORE: "Deleted code (Before)",
}

# Series labels for the four-bucket metrics, in stacking order (best risk to worst).
_RISK_SERIES_LABELS = ["Low risk", "Medium risk", "High risk", "Very high risk"]
# Duplication is reported as a binary distribution rather than four risk buckets.
_DUPLICATION_SERIES_LABELS = ["Non-redundant", "Redundant"]

# Risk distributions are normalized to sum to 100, so bars span the full axis.
_AXIS_MAX = 100
_EMPTY_RISK_BUCKETS = [0.0, 0.0, 0.0, 0.0]


def _series_labels(metric: MaintMetric) -> list[str]:
    if metric == MaintMetric.DUPLICATION:
        return _DUPLICATION_SERIES_LABELS
    return _RISK_SERIES_LABELS


def _column_series_values(
    metric: MaintMetric, risk_buckets: list[float]
) -> list[float]:
    """Turn a column's [low, moderate, high, very-high] buckets into series values.

    For duplication the four buckets collapse to non-redundant (low risk) and redundant
    (everything else); other metrics keep their four risk buckets.
    """
    normalized = normalize_percentages(risk_buckets)
    if metric == MaintMetric.DUPLICATION:
        non_redundant = normalized[0]
        redundant = sum(normalized[1:])
        return [non_redundant, redundant]
    return normalized


def _build_delta_chart_data(
    delta_type: DeltaType, metric: MaintMetric
) -> CategoryChartData:
    system_data = maintainability_delta_quality_system_by_type[delta_type]
    columns = system_data.columns

    values_per_column = [
        _column_series_values(
            metric, system_data.risk_buckets(metric, column) or _EMPTY_RISK_BUCKETS
        )
        for column in columns
    ]

    chart_data = CategoryChartData()
    chart_data.categories = [_COLUMN_LABELS[column] for column in columns]
    for series_index, label in enumerate(_series_labels(metric)):
        chart_data.add_series(
            label, [column_values[series_index] for column_values in values_per_column]
        )
    return chart_data


def _resolve_delta_chart(presentation: Presentation, key: str, value_cb: Callable):
    charts = rendering.pptx.find_charts(presentation, key)
    if not charts:
        return

    chart_data = value_cb()
    for chart in charts:
        chart.replace_data(chart_data)
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = _AXIS_MAX


class MaintainabilityDeltaQualityChartPlaceholder(ParameterizedPlaceholder):
    """Stacked bar chart of a maintainability metric's risk distribution, comparing the
    delta-quality code categories (new, changed or removed code against total code)."""

    __doc_type__ = PlaceholderDocType.CHART
    key = "DELTA_QUALITY_{type}_{metric}_CHART"
    allowed_parameters = MultiParameterList(list(DeltaType), DELTA_QUALITY_METRICS)

    @classmethod
    def value(cls, delta_type: DeltaType, metric: MaintMetric) -> CategoryChartData:
        return _build_delta_chart_data(delta_type, metric)

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        _resolve_delta_chart(presentation, key, value_cb)
