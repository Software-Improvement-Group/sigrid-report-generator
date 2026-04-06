#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from pptx.chart.data import XyChartData
from pptx.presentation import Presentation

from report_generator.generator.domain import (
    maintainability_data,
    maintainability_portfolio_data,
    system_metadata,
)
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.implementations.base import (
    Placeholder,
    PlaceholderDocType,
)


def _build_xy_chart_data(
    series_name: str, points: list[tuple[float, float]]
) -> XyChartData:
    # Volume values below 0.1 will not be displayed on the chart
    chart_data = XyChartData()
    series = chart_data.add_series(series_name)
    for x, y in points:
        series.add_data_point(max(x, 0.1), y)
    return chart_data


def _display_name(system_name: str) -> str:
    md = maintainability_portfolio_data.get_system_metadata(system_name)
    return md.get("displayName") or system_name


def _build_system_chart_data() -> XyChartData:
    return _build_xy_chart_data(
        system_metadata.display_name,
        [(maintainability_data.system_py, maintainability_data.maintainability_rating)],
    )


def _build_portfolio_chart_data() -> XyChartData:
    def to_point(system_name):
        snapshot = maintainability_portfolio_data.end_snapshot(system_name)
        return snapshot["volumeInPersonMonths"] / 12.0, snapshot["maintainability"]

    return _build_xy_chart_data(
        "Portfolio",
        [to_point(n) for n in maintainability_portfolio_data.system_names],
    )


def _populate_system_chart(presentation: Presentation, key: str, build_fn) -> None:
    charts = rendering.pptx.find_charts(presentation, key)
    if not charts:
        return
    chart_data = build_fn()
    for chart in charts:
        chart.replace_data(chart_data)
        for series in chart.series:
            series.points[0].data_label.text_frame.text = series.name


def _populate_portfolio_chart(presentation: Presentation) -> None:
    charts = rendering.pptx.find_charts(presentation, "PORTFOLIO_GALAXY_CHART")
    if not charts:
        return
    chart_data = _build_portfolio_chart_data()
    display_names = [
        _display_name(n) for n in maintainability_portfolio_data.system_names
    ]
    for chart in charts:
        chart.replace_data(chart_data)
        for i, point in enumerate(chart.series[0].points):
            point.data_label.text_frame.text = display_names[i]


class MaintainabilityGalaxyChartPlaceholder(Placeholder):
    """Traditional SIG benchmark galaxy chart."""

    key = "GALAXY_CHART"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls):
        return _build_system_chart_data()

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, _) -> None:
        _populate_system_chart(presentation, key, _build_system_chart_data)


class MaintainabilityPortfolioGalaxyChartPlaceholder(Placeholder):
    """Portfolio-level galaxy chart with one data point per system."""

    key = "PORTFOLIO_GALAXY_CHART"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls):
        return _build_portfolio_chart_data()

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, _) -> None:
        _populate_portfolio_chart(presentation)
