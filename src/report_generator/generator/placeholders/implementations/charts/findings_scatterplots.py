#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import math
from dataclasses import dataclass

from pptx.chart.data import XyChartData
from pptx.presentation import Presentation

from report_generator.generator.domain import (
    npr_5333_functional_suitability_portfolio_data,
    reliability_ratings_portfolio_data,
    security_ratings_portfolio_data,
)
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.implementations.base import (
    Placeholder,
    PlaceholderDocType,
)
from report_generator.generator.placeholders.implementations.charts.base import (
    findings_x_axis_max,
)


@dataclass(frozen=True)
class _FunctionalSuitabilityBounds:
    max_x: int
    max_y: float


def _functional_suitability_y_axis_max(max_ratio: float) -> float:
    return max(1.0, math.ceil(max_ratio * 10) / 10)


def _build_scatterplot_data(
    domain_data, series_name: str
) -> tuple[XyChartData, list[str], int]:
    findings_index = {
        entry["systemName"]: entry["findings_above_objective"]
        for entry in domain_data.findings_above_objective
    }
    points, display_names = _collect_points(domain_data, findings_index)
    chart_data = XyChartData()
    series = chart_data.add_series(series_name)
    for x, y in points:
        series.add_data_point(x, y)
    max_findings = findings_x_axis_max(max(findings_index.values(), default=0))
    return chart_data, display_names, max_findings


def _collect_points(
    domain_data, findings_by_name: dict[str, int]
) -> tuple[list, list[str]]:
    points = []
    display_names = []
    for system_name in domain_data.system_names:
        system = domain_data.get_system(system_name)
        if system is None:
            continue
        rating = system.get("rating")
        if rating is None:
            continue
        points.append((findings_by_name.get(system_name, 0), rating))
        display_names.append(domain_data.get_display_name(system_name))
    return points, display_names


def _populate_charts(
    charts, chart_data, display_names: list[str], max_findings: int
) -> None:
    for chart in charts:
        chart.replace_data(chart_data)
        chart.category_axis.minimum_scale = 0
        chart.category_axis.maximum_scale = max_findings
        for i, point in enumerate(chart.series[0].points):
            point.data_label.text_frame.text = display_names[i]


def _populate_functional_suitability_charts(
    charts, chart_data, display_names: list[str], bounds: _FunctionalSuitabilityBounds
) -> None:
    for chart in charts:
        chart.replace_data(chart_data)
        chart.category_axis.minimum_scale = 0
        chart.category_axis.maximum_scale = bounds.max_x
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = bounds.max_y
        for i, point in enumerate(chart.series[0].points):
            point.data_label.text_frame.text = display_names[i]


def _build_functional_suitability_chart_series(
    entries: list[dict],
) -> tuple[XyChartData, list[str]]:
    chart_data = XyChartData()
    series = chart_data.add_series("Functional Suitability")
    for entry in entries:
        series.add_data_point(entry["finding_count"], entry["test_code_ratio"])
    return chart_data, [entry["display_name"] for entry in entries]


def _build_functional_suitability_scatterplot_data() -> tuple[
    XyChartData, list[str], _FunctionalSuitabilityBounds
]:
    entries = [
        e
        for e in npr_5333_functional_suitability_portfolio_data.top_systems_by_finding_count
        if e["test_code_ratio"] is not None
    ]
    chart_data, display_names = _build_functional_suitability_chart_series(entries)
    max_x = findings_x_axis_max(max((e["finding_count"] for e in entries), default=0))
    max_y = _functional_suitability_y_axis_max(
        max((e["test_code_ratio"] for e in entries), default=0.0)
    )
    return chart_data, display_names, _FunctionalSuitabilityBounds(max_x, max_y)


def _resolve_scatterplot_pptx(
    presentation, key: str, domain_data, series_name: str
) -> None:
    charts = rendering.pptx.find_charts(presentation, key)
    if not charts:
        return
    chart_data, display_names, max_findings = _build_scatterplot_data(
        domain_data, series_name
    )
    _populate_charts(charts, chart_data, display_names, max_findings)


class PortfolioSecurityScatterplotPlaceholder(Placeholder):
    """Portfolio scatterplot: open security findings above objective (X) vs security rating (Y)."""

    key = "PORTFOLIO_SECURITY_SCATTERPLOT"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls):
        chart_data, _, __ = _build_scatterplot_data(
            security_ratings_portfolio_data, "Security"
        )
        return chart_data

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, _) -> None:
        _resolve_scatterplot_pptx(
            presentation, key, security_ratings_portfolio_data, "Security"
        )


class PortfolioReliabilityScatterplotPlaceholder(Placeholder):
    """Portfolio scatterplot: open reliability findings above objective (X) vs reliability rating (Y)."""

    key = "PORTFOLIO_RELIABILITY_SCATTERPLOT"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls):
        chart_data, _, __ = _build_scatterplot_data(
            reliability_ratings_portfolio_data, "Reliability"
        )
        return chart_data

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, _) -> None:
        _resolve_scatterplot_pptx(
            presentation, key, reliability_ratings_portfolio_data, "Reliability"
        )


class PortfolioNpr5333FunctionalSuitabilityScatterplotPlaceholder(Placeholder):
    """Portfolio scatterplot: NPR-5333 functional suitability findings (X) vs test code ratio (Y)."""

    key = "PORTFOLIO_NPR_5333_FUNCTIONAL_SUITABILITY_SCATTERPLOT"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls):
        chart_data, _, __ = _build_functional_suitability_scatterplot_data()
        return chart_data

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, _) -> None:
        charts = rendering.pptx.find_charts(presentation, key)
        if not charts:
            return
        chart_data, display_names, bounds = (
            _build_functional_suitability_scatterplot_data()
        )
        _populate_functional_suitability_charts(
            charts, chart_data, display_names, bounds
        )
