#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from datetime import date
from typing import Callable

from pptx.presentation import Presentation
from pptx.util import Pt

from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.implementations.base import Placeholder
from report_generator.generator.placeholders.rendering.common import FontProperties

_FONT_SIZE_TWO_DIGITS = 12
_FONT_SIZE_THREE_DIGITS = 10
_FONT_SIZE_FOUR_DIGITS = 7


def _font_size_for_days(days: int) -> int:
    if days < 100:
        return _FONT_SIZE_TWO_DIGITS
    elif days < 1000:
        return _FONT_SIZE_THREE_DIGITS
    else:
        return _FONT_SIZE_FOUR_DIGITS


class DaysSinceMythosDisclosurePlaceholder(Placeholder):
    """Number of days elapsed since Anthropic's initial Mythos disclosure on April 7, 2026."""

    key = "DAYS_SINCE_MYTHOS_DISCLOSURE"

    @classmethod
    def value(cls) -> str:
        return str((date.today() - date(2026, 4, 7)).days)

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        paragraphs = rendering.pptx.find_text_in_presentation(presentation, key)
        if not paragraphs:
            return
        value = value_cb()
        days = int(value)
        font = FontProperties(size=Pt(_font_size_for_days(days)))
        rendering.pptx.update_many_paragraphs(paragraphs, key, value, font)
