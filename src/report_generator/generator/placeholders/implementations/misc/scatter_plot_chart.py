#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from pptx.chart.data import BubbleChartData
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation

from report_generator.generator.domain import modernization_data
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.implementations.base import (
    Placeholder,
    PlaceholderDocType,
)

BLUE_GRADIENT = ["003DAB", "2E6BFF", "8DA8FF", "DBE1FF", "8A98A8"]


def _group_candidates() -> dict[str, list]:
    grouped = {c: [] for c in ["CRITICAL", "HIGH", "MEDIUM", "LOW", "UNKNOWN"]}
    for candidate in modernization_data.modernization_candidates:
        grouped[candidate.business_criticality.upper()].append(candidate)
    return grouped


def _build_chart_data(grouped_candidates: dict) -> BubbleChartData:
    chart_data = BubbleChartData()
    for group, candidates in grouped_candidates.items():
        series = chart_data.add_series(group.title())
        for c in candidates:
            series.add_data_point(
                c.estimated_effort_py, c.estimated_change_speed, c.volume_in_py
            )
    return chart_data


def _style_chart(chart, grouped_candidates: dict) -> None:
    for i, group in enumerate(grouped_candidates):
        for j, candidate in enumerate(grouped_candidates[group]):
            chart.series[i].points[
                j
            ].data_label.text_frame.text = candidate.display_name
            chart.series[i].points[j].format.line.color.rgb = RGBColor(255, 255, 255)
            chart.series[i].points[j].format.fill.solid()
            chart.series[i].points[j].format.fill.fore_color.rgb = RGBColor.from_string(
                BLUE_GRADIENT[i]
            )
    chart.value_axis.minimum_scale = 0
    chart.category_axis.minimum_scale = 0


def _populate_chart(presentation: Presentation, key: str) -> None:
    charts = rendering.pptx.find_charts(presentation, key)
    if not charts:
        return
    grouped_candidates = _group_candidates()
    chart_data = _build_chart_data(grouped_candidates)
    for chart in charts:
        chart.replace_data(chart_data)
        _style_chart(chart, grouped_candidates)


class ModernizationScatterPlotChartPlaceholder(Placeholder):
    key = "MODERNIZATION_SCATTER_PLOT_CHART"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _build_chart_data(_group_candidates())

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, _) -> None:
        _populate_chart(presentation, key)
