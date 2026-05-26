#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import statistics
from abc import ABC, abstractmethod
from typing import Callable

from pptx.presentation import Presentation

from report_generator.generator.domain import (
    osh_data,
    osh_portfolio_data,
    security_data,
    security_findings_portfolio_data,
)
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.formatting.urgency import (
    UrgencyColors,
    exploit_probability_colors,
    exploit_probability_label,
    library_age_colors,
    library_age_label,
    urgency_colors,
    urgency_width,
)
from report_generator.generator.placeholders.implementations.base import Placeholder
from report_generator.generator.placeholders.rendering.common import (
    FontColor,
    FontProperties,
)
from report_generator.generator.placeholders.rendering.pptx import ShapeProperties


def _apply_colored_shape(
    shapes,
    paragraphs,
    key: str,
    shape_props: ShapeProperties,
    display_value: str,
    text_color=None,
):
    for shape in shapes:
        rendering.pptx.apply_shape_properties(shape, shape_props)
    font = (
        FontProperties(color=FontColor(rgb=text_color))
        if text_color is not None
        else None
    )
    rendering.pptx.update_many_paragraphs(paragraphs, key, display_value, font)


class AbstractUrgencyShapePlaceholder(Placeholder, ABC):
    """Colors a shape and sets text color based on urgency, with a right-anchored width."""

    @classmethod
    @abstractmethod
    def _get_colors(cls) -> UrgencyColors: ...

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        shapes = rendering.pptx.find_shapes_with_text(presentation, key)
        paragraphs = rendering.pptx.find_text_in_presentation(presentation, key)
        if not shapes and not paragraphs:
            return
        colors = cls._get_colors()
        display_value = value_cb()
        _apply_colored_shape(
            shapes,
            paragraphs,
            key,
            shape_props=ShapeProperties(
                color=colors.shape,
                width_inches=urgency_width(display_value),
                width_anchor_right=True,
            ),
            display_value=display_value,
            text_color=colors.text,
        )


class SecurityCVSSUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape and sets text based on the urgency of security findings."""

    key = "SECURITY_CVSS_URGENCY"

    @classmethod
    def value(cls):
        return "review"

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return urgency_colors(
            {
                "critical": security_data.count_findings("CRITICAL"),
                "high": security_data.count_findings("HIGH"),
                "medium": security_data.count_findings("MEDIUM"),
                "low": security_data.count_findings("LOW"),
            }
        )


class SecurityPortfolioCVSSUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape and sets text based on the urgency of security findings across the portfolio."""

    key = "SECURITY_PORTFOLIO_CVSS_URGENCY"

    @classmethod
    def value(cls):
        return "review"

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return urgency_colors(
            {
                "critical": security_findings_portfolio_data.count_findings("CRITICAL"),
                "high": security_findings_portfolio_data.count_findings("HIGH"),
                "medium": security_findings_portfolio_data.count_findings("MEDIUM"),
                "low": security_findings_portfolio_data.count_findings("LOW"),
            }
        )


class OSHKnownVulnerabilitiesUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape red, yellow, or green based on the urgency of known vulnerabilities."""

    key = "OSH_KNOWN_VULNERABILITIES_URGENCY"

    @classmethod
    def value(cls):
        return "review"

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return urgency_colors(osh_data.vulnerability_distribution)


class OSHAverageLibraryAgeUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape red, orange, yellow, or green based on the average library age."""

    key = "OSH_AVERAGE_LIBRARY_AGE_URGENCY"

    @classmethod
    def value(cls):
        return library_age_label(statistics.mean(osh_data.age_distribution))

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return library_age_colors(statistics.mean(osh_data.age_distribution))


class OSHExploitProbabilityUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape red, orange, yellow, or green based on the probability of exploit."""

    key = "OSH_PROBABILITY_OF_EXPLOIT_URGENCY"

    @classmethod
    def value(cls):
        return exploit_probability_label(osh_data.exploit_probability)

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return exploit_probability_colors(osh_data.exploit_probability)


class OSHPortfolioKnownVulnerabilitiesUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape red, yellow, or green based on the urgency of known vulnerabilities across the portfolio."""

    key = "OSH_PORTFOLIO_KNOWN_VULNERABILITIES_URGENCY"

    @classmethod
    def value(cls):
        return "review"

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return urgency_colors(osh_portfolio_data.vulnerability_distribution)


class OSHPortfolioAverageLibraryAgeUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape red, orange, yellow, or green based on the average library age across the portfolio."""

    key = "OSH_PORTFOLIO_AVERAGE_LIBRARY_AGE_URGENCY"

    @classmethod
    def value(cls):
        return library_age_label(statistics.mean(osh_portfolio_data.age_distribution))

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return library_age_colors(statistics.mean(osh_portfolio_data.age_distribution))


class OSHPortfolioExploitProbabilityUrgency(AbstractUrgencyShapePlaceholder):
    """Colors a shape red, orange, yellow, or green based on the probability of exploit across the portfolio."""

    key = "OSH_PORTFOLIO_PROBABILITY_OF_EXPLOIT_URGENCY"

    @classmethod
    def value(cls):
        return exploit_probability_label(osh_portfolio_data.exploit_probability)

    @classmethod
    def _get_colors(cls) -> UrgencyColors:
        return exploit_probability_colors(osh_portfolio_data.exploit_probability)
