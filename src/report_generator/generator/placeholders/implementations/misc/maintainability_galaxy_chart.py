#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from pptx.chart.data import XyChartData
from pptx.presentation import Presentation

from report_generator.generator.domain import maintainability_data, system_metadata
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.implementations.base import (
    Placeholder,
    PlaceholderDocType,
)


def _build_chart_data() -> XyChartData:
    chart_data = XyChartData()
    series = chart_data.add_series("Series 1")
    # Correct volume to be at least 0.1, anything lower will not be displayed on the chart
    series.add_data_point(
        max(maintainability_data.system_py, 0.1),
        maintainability_data.maintainability_rating,
    )
    return chart_data


def _populate_chart(presentation: Presentation, key: str) -> None:
    charts = rendering.pptx.find_charts(presentation, key)
    if not charts:
        return
    chart_data = _build_chart_data()
    system_name = system_metadata.display_name
    for chart in charts:
        chart.replace_data(chart_data)
        chart.series[0].points[0].data_label.text_frame.text = system_name


class MaintainabilityGalaxyChartPlaceholder(Placeholder):
    """Traditional SIG benchmark galaxy chart."""

    key = "GALAXY_SLIDE"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls):
        return _build_chart_data()

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, _) -> None:
        _populate_chart(presentation, key)
