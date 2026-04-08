#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from abc import ABC, abstractmethod
from typing import Callable

from pptx.chart.data import CategoryChartData
from pptx.presentation import Presentation

from report_generator.generator.domain import (
    maintainability_data,
    modernization_data,
    objectives_data,
    progress_sigrid_data,
    sigrid_hygiene_portfolio_data,
)
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.formatting import formatters
from report_generator.generator.placeholders.implementations.base import (
    Placeholder,
    PlaceholderDocType,
)


def _build_chart_data(values: dict) -> CategoryChartData:
    chart_data = CategoryChartData()
    chart_data.categories = values["labels"]
    series = values["series"]
    series_names = values.get("seriesNames", [])

    # Case 1: No custom series names -> use axisLabel for all series
    if not series_names:
        for y in values["series"]:
            chart_data.add_series(values["axisLabel"], y)
    # Case 2: Exact match between names and series -> map 1:1
    elif len(series_names) == len(series):
        for name, y in zip(series_names, series):
            chart_data.add_series(name, y)
    else:
        raise ValueError(
            f"seriesNames length ({len(series_names)}) does not match "
            f"series length ({len(series)})."
        )

    return chart_data


def _apply_colors(chart, colors: list) -> None:
    for serie in chart.series:
        for idx, point in enumerate(serie.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[idx]


def _populate_chart(presentation: Presentation, key: str, value_cb: Callable) -> None:
    charts = rendering.pptx.find_charts(presentation, key)
    if not charts:
        return
    values = value_cb()
    chart_data = _build_chart_data(values)
    colors = values["colors"]
    for chart in charts:
        chart.replace_data(chart_data)
        if colors:
            _apply_colors(chart, colors)


class _AbstractCategoryChartPlaceholder(Placeholder, ABC):
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    @abstractmethod
    def labels(cls):
        pass

    @classmethod
    @abstractmethod
    def series(cls):
        pass

    @classmethod
    def series_names(cls):
        return []

    @classmethod
    def colors(cls):
        return []

    @classmethod
    @abstractmethod
    def axis_label(cls):
        pass

    @classmethod
    def value(cls):
        return {
            "labels": cls.labels(),
            "series": cls.series(),
            "seriesNames": cls.series_names(),
            "colors": cls.colors(),
            "axisLabel": cls.axis_label(),
        }

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable):
        _populate_chart(presentation, key, value_cb)


class TechnologyCategoryChartPlaceholder(_AbstractCategoryChartPlaceholder):
    """Chart with volume (in % of person months of code) per technology."""

    key = "TECHNOLOGY_CHART"

    @classmethod
    def labels(cls):
        return [data["displayName"] for data in maintainability_data.sorted_tech]

    @classmethod
    def series(cls):
        return [
            [
                data["volumeInPersonMonths"] / maintainability_data.tech_total_volume_pm
                for data in maintainability_data.sorted_tech
            ]
        ]

    @classmethod
    def axis_label(cls):
        return "Volume in Person Months"


class TestCodeRatioCategoryChartPlaceholder(_AbstractCategoryChartPlaceholder):
    """Pie chart with volume and % of test code per technology, colored in line with the SIG test code benchmark."""

    key = "TEST_CODE_RATIO_CHART"

    @classmethod
    def labels(cls):
        return [data["displayName"] for data in maintainability_data.sorted_tech]

    @classmethod
    def series(cls):
        return [
            [
                data["volumeInPersonMonths"] / maintainability_data.tech_total_volume_pm
                for data in maintainability_data.sorted_tech
            ]
        ]

    @classmethod
    def colors(cls):
        return [
            rendering.pptx.test_code_ratio_color(data["testCodeRatio"])
            for data in maintainability_data.sorted_tech
        ]

    @classmethod
    def axis_label(cls):
        return "Volume in Person Months"


class TechnicalDebtSystemsChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "TECHNICAL_DEBT_SYSTEMS_CHART"

    @classmethod
    def labels(cls):
        candidates = modernization_data.modernization_candidates_by_estimated_effort[
            0:20
        ]
        return [candidate.display_name for candidate in candidates]

    @classmethod
    def series(cls):
        candidates = modernization_data.modernization_candidates_by_estimated_effort[
            0:20
        ]
        technical_debt = [candidate.estimated_effort_py for candidate in candidates]
        remaining = [
            candidate.volume_in_py - candidate.estimated_effort_py
            for candidate in candidates
        ]
        return [technical_debt, remaining]

    @classmethod
    def axis_label(cls):
        return "Code volume in person years"


class ObjectivesOverallChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_OVERALL_CHART"

    @classmethod
    def labels(cls):
        return [period.start.strftime("%m/%Y") for period in objectives_data.periods]

    @classmethod
    def series(cls):
        return objectives_data.get_portfolio_trend_series(None)

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesMaintainabilityChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_MAINTAINABILITY_CHART"

    @classmethod
    def labels(cls):
        return [period.start.strftime("%m/%Y") for period in objectives_data.periods]

    @classmethod
    def series(cls):
        return objectives_data.get_portfolio_trend_series("MAINTAINABILITY")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesArchitectureChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_ARCHITECTURE_CHART"

    @classmethod
    def labels(cls):
        return [period.start.strftime("%m/%Y") for period in objectives_data.periods]

    @classmethod
    def series(cls):
        return objectives_data.get_portfolio_trend_series("ARCHITECTURE_QUALITY")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesSecurityChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_SECURITY_CHART"

    @classmethod
    def labels(cls):
        return [period.start.strftime("%m/%Y") for period in objectives_data.periods]

    @classmethod
    def series(cls):
        return objectives_data.get_portfolio_trend_series("SECURITY")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesOpenSourceHealthChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_OSH_CHART"

    @classmethod
    def labels(cls):
        return [period.start.strftime("%m/%Y") for period in objectives_data.periods]

    @classmethod
    def series(cls):
        return objectives_data.get_portfolio_trend_series("OPEN_SOURCE_HEALTH")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesStatusChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_STATUS_CHART"

    @classmethod
    def labels(cls):
        return ["Overall"]

    @classmethod
    def series(cls):
        return objectives_data.get_portfolio_status_series()

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesTeamChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_TEAM_CHART"

    @classmethod
    def labels(cls):
        return list(objectives_data.teams.keys())

    @classmethod
    def series(cls):
        return objectives_data.get_team_status_series()

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesCapabilitiesChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "OBJECTIVES_CAPABILITY_CHART"

    @classmethod
    def labels(cls):
        return [
            capability.title().replace("_", " ")
            for capability in objectives_data.capabilities
        ]

    @classmethod
    def series(cls):
        return objectives_data.get_capability_status_series()

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesOverallChartSigridPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "PROGRESS_TIME_CHART"

    @classmethod
    def labels(cls):
        return [
            period.start.strftime("%m/%Y") for period in progress_sigrid_data.periods
        ]

    @classmethod
    def series(cls):
        return progress_sigrid_data.get_portfolio_trend_series(None)

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesMaintainabilityChartSigridPlaceholder(
    _AbstractCategoryChartPlaceholder
):
    key = "PROGRESS_MAINTAINABILITY_TIME_CHART"

    @classmethod
    def labels(cls):
        return [
            period.start.strftime("%m/%Y") for period in progress_sigrid_data.periods
        ]

    @classmethod
    def series(cls):
        return progress_sigrid_data.get_portfolio_trend_series("MAINTAINABILITY")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesArchitectureChartSigridPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "PROGRESS_ARCHITECTURE_TIME_CHART"

    @classmethod
    def labels(cls):
        return [
            period.start.strftime("%m/%Y") for period in progress_sigrid_data.periods
        ]

    @classmethod
    def series(cls):
        return progress_sigrid_data.get_portfolio_trend_series("ARCHITECTURE_QUALITY")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesSecurityChartSigridPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "PROGRESS_SECURITY_TIME_CHART"

    @classmethod
    def labels(cls):
        return [
            period.start.strftime("%m/%Y") for period in progress_sigrid_data.periods
        ]

    @classmethod
    def series(cls):
        return progress_sigrid_data.get_portfolio_trend_series("SECURITY")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesOpenSourceHealthChartSigridPlaceholder(
    _AbstractCategoryChartPlaceholder
):
    key = "PROGRESS_OSH_TIME_CHART"

    @classmethod
    def labels(cls):
        return [
            period.start.strftime("%m/%Y") for period in progress_sigrid_data.periods
        ]

    @classmethod
    def series(cls):
        return progress_sigrid_data.get_portfolio_trend_series("OPEN_SOURCE_HEALTH")

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesCapabilitiesChartSigridPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "PROGRESS_CAPABILITY_CHART"

    @classmethod
    def labels(cls):
        return [
            capability.title().replace("_", " ")
            for capability in progress_sigrid_data.capabilities
        ]

    @classmethod
    def series(cls):
        return progress_sigrid_data.get_capability_status_series()

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class ObjectivesStatusChartSigridPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "PROGRESS_STATUS_CHART"

    @classmethod
    def labels(cls):
        return ["Overall"]

    @classmethod
    def series(cls):
        return progress_sigrid_data.get_portfolio_status_series()

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class MetadataCompletenessChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "METADATA_COMPLETENESS_CHART"

    @classmethod
    def labels(cls):
        return [
            formatters.from_json_name(key)
            for key in sigrid_hygiene_portfolio_data.metadata_completeness.keys()
        ]

    @classmethod
    def series(cls):
        values = list(sigrid_hygiene_portfolio_data.metadata_completeness.values())
        return [
            [v[0] for v in values],
            [v[1] for v in values],
        ]

    @classmethod
    def axis_label(cls):
        return "Percentage of portfolio"


class SnapshotFreshnessChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "SNAPSHOT_FRESHNESS_CHART"

    @classmethod
    def labels(cls):
        return ["Total", "1 week", "1 month", "3 months", "6 months", ">6 months"]

    @classmethod
    def series(cls):
        freshness_days = sigrid_hygiene_portfolio_data.snapshot_freshness
        result = [
            formatters.split_days_into_buckets(
                freshness_days.values(), buckets=[7, 30, 90, 180]
            )
        ]
        return result

    @classmethod
    def axis_label(cls):
        return "Systems"


class EolDeactivatedSystemsChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "EOL_DEACTIVATED_CHART"

    @classmethod
    def labels(cls):
        return ["Total", "Deactivated", "EOL", "EOL & Deactivated"]

    @classmethod
    def series(cls):
        return sigrid_hygiene_portfolio_data.eol_deactivated_systems

    @classmethod
    def axis_label(cls):
        return "Systems"


class UsersLastLoginChartPlaceholder(_AbstractCategoryChartPlaceholder):
    key = "USERS_LAST_LOGIN_CHART"

    @classmethod
    def labels(cls):
        return ["Total", "1 week", "1 month", "3 months", "1 year", ">1 year"]

    @classmethod
    def series(cls):
        buckets = [7, 30, 90, 365]
        days_admin = sigrid_hygiene_portfolio_data.get_last_access_time_users(
            role="ADMIN"
        )
        days_maintainer = sigrid_hygiene_portfolio_data.get_last_access_time_users(
            role="MAINTAINER"
        )
        days_user = sigrid_hygiene_portfolio_data.get_last_access_time_users(
            role="USER"
        )
        return [
            formatters.split_days_into_buckets(days_admin, buckets=buckets),
            formatters.split_days_into_buckets(days_maintainer, buckets=buckets),
            formatters.split_days_into_buckets(days_user, buckets=buckets),
        ]

    @classmethod
    def series_names(cls):
        return ["Admin", "Maintainer", "User"]

    @classmethod
    def axis_label(cls):
        return "Users"
