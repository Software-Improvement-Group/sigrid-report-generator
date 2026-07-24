#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from abc import ABC
from collections.abc import Callable
from typing import ClassVar

from pptx.presentation import Presentation

from report_generator.generator.domain import architecture_data, maintainability_data
from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.formatting import formatters
from report_generator.generator.placeholders.implementations.base import (
    ParameterizedPlaceholder,
)
from report_generator.generator.placeholders.rendering.common import (
    FontColor,
    FontProperties,
)
from report_generator.generator.placeholders.rendering.pptx import ShapeProperties
from report_generator.generator.utils.constants import (
    ArchMetric,
    ArchSubcharacteristic,
    MaintMetric,
    MetricEnum,
)


def _apply_colored_shape(
    shapes,
    paragraphs,
    key: str,
    shape_props: ShapeProperties,
    display_value: str,
    text_color=None,
):
    for shape in shapes:
        rendering.pptx.apply_shape_properties(shape, shape_props)
    font = (
        FontProperties(color=FontColor(rgb=text_color))
        if text_color is not None
        else None
    )
    rendering.pptx.update_many_paragraphs(paragraphs, key, display_value, font)


class AbstractColorRatingPlaceholder(ParameterizedPlaceholder, ABC):
    """Fills this rating value and colors the shape to correspond to the rating color (e.g. yellow for 3 stars)."""

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        shapes = rendering.pptx.find_shapes_with_text(presentation, key)
        paragraphs = rendering.pptx.find_text_in_presentation(presentation, key)
        if not shapes and not paragraphs:
            return
        rating = value_cb()
        _apply_colored_shape(
            shapes,
            paragraphs,
            key,
            shape_props=ShapeProperties(
                color=rendering.pptx.determine_rating_color(rating)
            ),
            display_value=formatters.star_rating_round(rating),
        )


class MaintColorRatingPlaceholder(AbstractColorRatingPlaceholder):
    """Fills the rating value and colors the shape to the corresponding rating color for a maintainability metric."""

    key = "COLOR_MAINT_RATING_{parameter}"
    allowed_parameters: ClassVar[list] = list(MaintMetric)

    @classmethod
    def value(cls, metric: MaintMetric):
        metric_key = metric.to_json_name()
        return maintainability_data.data[metric_key]


class ArchColorRatingPlaceholder(AbstractColorRatingPlaceholder):
    """Fills the rating value and colors the shape to the corresponding rating color for an architecture metric."""

    key = "COLOR_ARCH_RATING_{parameter}"
    allowed_parameters: ClassVar[list] = list(ArchMetric) + list(ArchSubcharacteristic)

    @classmethod
    def value(cls, metric: MetricEnum):
        metric_key = metric.to_json_name()
        return architecture_data.get_score_for_prop_or_subchar(metric_key)
