#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from abc import ABC
from typing import Callable, Optional, Union

from docx.document import Document
from pptx.presentation import Presentation

from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.formatting import formatters
from report_generator.generator.placeholders.implementations.base import (
    MultiParameterList,
    ParameterizedPlaceholder,
    ParameterList,
    Placeholder,
    PlaceholderDocType,
    function_name_to_placeholder_key,
)
from report_generator.generator.placeholders.rendering.common import (
    FontColor,
    FontProperties,
)


class _DocumentAdapter:
    def __init__(self, find_func, update_func):
        self.find_text = find_func
        self.update_paragraphs = update_func


class _AbstractTextPlaceholder(Placeholder, ABC):
    __doc_type__ = PlaceholderDocType.TEXT

    _PPTX_ADAPTER = _DocumentAdapter(
        rendering.pptx.find_text_in_presentation, rendering.pptx.update_many_paragraphs
    )

    _DOCX_ADAPTER = _DocumentAdapter(
        rendering.docx.find_text_in_document, rendering.docx.update_many_paragraphs
    )

    @staticmethod
    def _resolve_with_adapter(
        adapter: _DocumentAdapter, document, key: str, value_cb: Callable[[], str]
    ) -> None:
        paragraphs = adapter.find_text(document, key)

        if len(paragraphs) == 0:
            return

        value = value_cb()
        if value is None:
            raise ValueError(f"Value for placeholder '{key}' is None")

        adapter.update_paragraphs(paragraphs, key, value)

    @staticmethod
    def resolve_pptx(
        presentation: Presentation, key: str, value_cb: Callable[[], str]
    ) -> None:
        _AbstractTextPlaceholder._resolve_with_adapter(
            _AbstractTextPlaceholder._PPTX_ADAPTER, presentation, key, value_cb
        )

    @staticmethod
    def resolve_docx(document: Document, key: str, value_cb: Callable[[], str]) -> None:
        _AbstractTextPlaceholder._resolve_with_adapter(
            _AbstractTextPlaceholder._DOCX_ADAPTER, document, key, value_cb
        )


def text_placeholder(
    custom_key: Optional[str] = None,
) -> Callable[[Callable[[], str]], type[Placeholder]]:
    def decorator(value_func: Callable[[], str]) -> type[Placeholder]:
        class TextPlaceholder(_AbstractTextPlaceholder):
            __doc__ = value_func.__doc__ if value_func.__doc__ else None
            key = (
                custom_key
                if custom_key
                else function_name_to_placeholder_key(value_func.__name__)
            )

            @classmethod
            def value(cls) -> str:
                return value_func()

        return TextPlaceholder

    return decorator


def _render_colored_delta(
    presentation: Presentation, key: str, delta_func: Callable[[], float]
) -> None:
    paragraphs = rendering.pptx.find_text_in_presentation(presentation, key)
    if not paragraphs:
        return
    delta = delta_func()
    font = FontProperties(
        color=FontColor(
            rgb=rendering.pptx.sentiment_color(formatters.delta_sentiment(delta))
        )
    )
    text = formatters.format_signed_delta(delta)
    rendering.pptx.update_many_paragraphs(paragraphs, key, text, font)


def delta_text_placeholder(
    custom_key: Optional[str] = None,
) -> Callable[[Callable[[], float]], type[Placeholder]]:
    """Turn a function returning a numeric delta into a text placeholder that renders the signed
    delta (e.g. +0.01, -0.01, =) colored green for an increase, red for a decrease and blue when
    unchanged. Coloring is applied in PowerPoint; Word renders the value without color."""

    def decorator(delta_func: Callable[[], float]) -> type[Placeholder]:
        class DeltaTextPlaceholder(_AbstractTextPlaceholder):
            __doc__ = delta_func.__doc__ if delta_func.__doc__ else None
            key = (
                custom_key
                if custom_key
                else function_name_to_placeholder_key(delta_func.__name__)
            )

            @classmethod
            def value(cls) -> str:
                return formatters.format_signed_delta(delta_func())

            @classmethod
            def resolve_pptx(
                cls, presentation: Presentation, key: str, value_cb: Callable
            ) -> None:
                _render_colored_delta(presentation, key, delta_func)

        return DeltaTextPlaceholder

    return decorator


def _render_colored_market_average(
    presentation: Presentation, key: str, score_func: Callable[[], float]
) -> None:
    paragraphs = rendering.pptx.find_text_in_presentation(presentation, key)
    if not paragraphs:
        return
    score = score_func()
    font = FontProperties(
        color=FontColor(
            rgb=rendering.pptx.sentiment_color(
                formatters.market_average_sentiment(score)
            )
        )
    )
    text = formatters.format_market_average(score)
    rendering.pptx.update_many_paragraphs(paragraphs, key, text, font)


def market_average_text_placeholder(
    custom_key: Optional[str] = None,
) -> Callable[[Callable[[], float]], type[Placeholder]]:
    """Turn a function returning a star rating into a text placeholder that renders whether the
    score is at market average: 'below' colored red (< 2.5), 'average' colored blue (2.5 - 3.4)
    and 'above' colored green (>= 3.5). Coloring is applied in PowerPoint; Word renders the value
    without color."""

    def decorator(score_func: Callable[[], float]) -> type[Placeholder]:
        class MarketAverageTextPlaceholder(_AbstractTextPlaceholder):
            __doc__ = score_func.__doc__ if score_func.__doc__ else None
            key = (
                custom_key
                if custom_key
                else function_name_to_placeholder_key(score_func.__name__)
            )

            @classmethod
            def value(cls) -> str:
                return formatters.format_market_average(score_func())

            @classmethod
            def resolve_pptx(
                cls, presentation: Presentation, key: str, value_cb: Callable
            ) -> None:
                _render_colored_market_average(presentation, key, score_func)

        return MarketAverageTextPlaceholder

    return decorator


def parameterized_text_placeholder(
    custom_key: str, parameters: Union[ParameterList, MultiParameterList]
) -> Callable:
    def decorator(value_func) -> type[ParameterizedPlaceholder]:
        class ParameterizedTextPlaceholder(
            ParameterizedPlaceholder, _AbstractTextPlaceholder
        ):
            __doc__ = value_func.__doc__ if value_func.__doc__ else None
            key = custom_key
            allowed_parameters = parameters

            @classmethod
            def value(cls, *args) -> str:
                return value_func(*args)

        return ParameterizedTextPlaceholder

    return decorator
