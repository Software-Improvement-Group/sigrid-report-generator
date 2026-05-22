#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from abc import ABC, abstractmethod
from typing import Callable

from pptx.presentation import Presentation

from report_generator.generator.placeholders import rendering
from report_generator.generator.placeholders.formatting import formatters
from report_generator.generator.placeholders.implementations.base import (
    ParameterizedPlaceholder,
    Placeholder,
)
from report_generator.generator.placeholders.implementations.text.shared.urgency import (
    UrgencyColors,
    urgency_width,
)
from report_generator.generator.placeholders.rendering.common import (
    FontColor,
    FontProperties,
)
from report_generator.generator.placeholders.rendering.pptx import ShapeProperties


def _apply_colored_shape(
    shapes,
    paragraphs,
    key: str,
    shape_props: ShapeProperties,
    display_value: str,
    text_color=None,
):
    for shape in shapes:
        rendering.pptx.apply_shape_properties(shape, shape_props)
    font = (
        FontProperties(color=FontColor(rgb=text_color))
        if text_color is not None
        else None
    )
    rendering.pptx.update_many_paragraphs(paragraphs, key, display_value, font)


class AbstractUrgencyShapePlaceholder(Placeholder, ABC):
    """Colors a shape and sets text color based on urgency, with a right-anchored width."""

    @classmethod
    @abstractmethod
    def _get_colors(cls) -> UrgencyColors: ...

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        shapes = rendering.pptx.find_shapes_with_text(presentation, key)
        paragraphs = rendering.pptx.find_text_in_presentation(presentation, key)
        if not shapes and not paragraphs:
            return
        colors = cls._get_colors()
        display_value = value_cb()
        _apply_colored_shape(
            shapes,
            paragraphs,
            key,
            shape_props=ShapeProperties(
                color=colors.shape,
                width_inches=urgency_width(display_value),
                width_anchor_right=True,
            ),
            display_value=display_value,
            text_color=colors.text,
        )


class AbstractColorRatingPlaceholder(ParameterizedPlaceholder, ABC):
    """Fills this rating value and colors the shape to correspond to the rating color (e.g. yellow for 3 stars)."""

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        shapes = rendering.pptx.find_shapes_with_text(presentation, key)
        paragraphs = rendering.pptx.find_text_in_presentation(presentation, key)
        if not shapes and not paragraphs:
            return
        rating = value_cb()
        _apply_colored_shape(
            shapes,
            paragraphs,
            key,
            shape_props=ShapeProperties(
                color=rendering.pptx.determine_rating_color(rating)
            ),
            display_value=formatters.star_rating_round(rating),
        )
