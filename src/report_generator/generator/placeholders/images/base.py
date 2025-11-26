#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from typing import Callable
from abc import ABC

from pptx.presentation import Presentation

from report_generator.generator.placeholders.base import Placeholder, PlaceholderDocType
from report_generator.generator import report_utils

import io
from pptx.util import Inches
import logging

class _AbstractImagePlaceholder(Placeholder, ABC):
    __doc_type__ = PlaceholderDocType.IMAGE
    BUNDLE_COLOR = f"#{report_utils.pptx.SIG_GREY_COLOR}"
    NA_STAR_COLOR = f"#{report_utils.pptx.NA_STAR_COLOR}"

    @classmethod
    def resolve_pptx(cls, presentation: Presentation, key: str, value_cb: Callable):
        shapes = report_utils.pptx.find_shapes(presentation, key)
        if len(shapes) == 0:
            return

        for shape in shapes:
            fig = value_cb(parameter={'height':shape.height.inches, 'width':shape.width.inches})
            cls.create_and_add_image_to_slide(shape, fig)
    
    @staticmethod
    def create_and_add_image_to_slide(shape_placeholder, fig):
        pos_left = shape_placeholder.left.inches
        pos_top = shape_placeholder.top.inches
        pos_width = shape_placeholder.width.inches
        pos_height = shape_placeholder.height.inches

        if fig is None:
            logging.warning("Figure data of an image placeholder is None.")
            return

        with io.BytesIO() as buf:
            fig.savefig(buf, dpi='figure', bbox_inches='tight', transparent=True, pad_inches=0)
            buf.seek(0)
            
            shape_placeholder.part.slide.shapes.add_picture(buf,
                left=Inches(pos_left), top=Inches(pos_top),
                width=Inches(pos_width), height=Inches(pos_height))

        el = shape_placeholder.element
        el.getparent().remove(el)
