#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

from typing import Callable

from pptx.chart.data import CategoryChartData
from pptx.presentation import Presentation

from report_generator.generator import report_utils
from report_generator.generator.data_models import security_dashboard_findings_portfolio_data
from report_generator.generator.data_models import security_dashboard_resolution_times_portfolio_data
from report_generator.generator.placeholders import Placeholder
from report_generator.generator.placeholders.base import PlaceholderDocType


def _add_month_data_row(arrays_dict, month, new_val, existing_val, resolved_val, total_val):
    """Add a single row of data for a month"""
    arrays_dict['categories'].append(month)
    arrays_dict['new'].append(new_val)
    arrays_dict['existing'].append(existing_val)
    arrays_dict['resolved'].append(resolved_val)
    arrays_dict['total'].append(total_val)


def _build_chart_data_arrays(data):
    """Build data arrays for the grouped/clustered chart structure.
    
    For the chart to work with stacked (New+Existing) and clustered (Resolved) bars,
    we need to create categories with blank separators for grouping.
    
    Returns arrays in the format matching the Excel structure with grouped data.
    """
    columns = data['columns']
    new = data['new']
    existing = data['existing']
    resolved = data['resolved']
    
    # Initialize result dictionary
    arrays_dict = {
        'categories': [],
        'new': [],
        'existing': [],
        'resolved': [],
        'total': []
    }
    
    for i, month in enumerate(columns):
        # Calculate total for this month (new + existing)
        total = new[i] + existing[i]
        
        # Row 1: Month name with New and Existing stacked (no Resolved here)
        _add_month_data_row(arrays_dict, month, new[i], existing[i], 0, total)
        
        # Row 2: Blank for spacing (only Resolved gets value here for clustering effect)
        _add_month_data_row(arrays_dict, '', 0, 0, resolved[i], resolved[i])
        
        # Row 3: Blank for spacing (skip for last month)
        if i < len(columns) - 1:
            _add_month_data_row(arrays_dict, '', 0, 0, 0, 0)
    
    return arrays_dict


def _create_security_findings_chart_data(severity):
    """Create CategoryChartData for security findings chart"""
    # Get aggregated data from data model
    aggregated_data = security_dashboard_findings_portfolio_data.chart_findings_by_severity(severity)
    
    # Build chart-ready arrays with grouping
    chart_arrays = _build_chart_data_arrays(aggregated_data)
    
    # Create CategoryChartData
    chart_data = CategoryChartData()
    chart_data.categories = chart_arrays['categories']
    
    # Add series: New and Existing will be stacked, Resolved will be clustered, Total displays on top
    chart_data.add_series('New', chart_arrays['new'])
    chart_data.add_series('Existing', chart_arrays['existing'])
    chart_data.add_series('Resolved', chart_arrays['resolved'])
    chart_data.add_series('Total', chart_arrays['total'])
    
    return chart_data


def _populate_chart(presentation: Presentation, key: str, value_cb: Callable, chart_name: str):
    """Populate a chart in the presentation by finding slides with the key and updating the named chart."""
    slides = report_utils.pptx.identify_specific_slide(presentation, key)
    
    if len(slides) == 0:
        return
    
    # Get chart data
    chart_data = value_cb()
    
    # Find and update charts by name
    for slide in slides:
        shapes_by_name = dict((s.name, s) for s in slide.shapes)
        if chart_name in shapes_by_name:
            chart = shapes_by_name[chart_name].chart
            chart.replace_data(chart_data)


class SecurityDashboardCriticalFindingsChartPlaceholder(Placeholder):
    """PowerPoint chart depicting new, existing, and resolved critical security findings over the last 12 months.
    
    Uses a combination of stacked and clustered columns where New and Existing are stacked together,
    and Resolved appears as a separate clustered column next to them for each month.
    """
    
    key = "PORTFOLIO_SECURITY_FINDINGS_CRITICAL"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _create_security_findings_chart_data("CRITICAL")

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        _populate_chart(presentation, key, value_cb, "PORTFOLIO_SECURITY_FINDINGS_CRITICAL")


class SecurityDashboardHighFindingsChartPlaceholder(Placeholder):
    """PowerPoint chart depicting new, existing, and resolved high security findings over the last 12 months.
    
    Uses a combination of stacked and clustered columns where New and Existing are stacked together,
    and Resolved appears as a separate clustered column next to them for each month.
    """
    
    key = "PORTFOLIO_SECURITY_FINDINGS_HIGHMED"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _create_security_findings_chart_data("HIGH")

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        _populate_chart(presentation, key, value_cb, "PORTFOLIO_SECURITY_FINDINGS_HIGH")


class SecurityDashboardMediumFindingsChartPlaceholder(Placeholder):
    """PowerPoint chart depicting new, existing, and resolved medium security findings over the last 12 months.
    
    Uses a combination of stacked and clustered columns where New and Existing are stacked together,
    and Resolved appears as a separate clustered column next to them for each month.
    """
    
    key = "PORTFOLIO_SECURITY_FINDINGS_HIGHMED"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _create_security_findings_chart_data("MEDIUM")

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        _populate_chart(presentation, key, value_cb, "PORTFOLIO_SECURITY_FINDINGS_MEDIUM")


# Resolution Times Placeholders

def _build_resolution_times_chart_data_arrays(data):
    """Build data arrays for stacked resolution times chart.
    
    Returns arrays with month categories and stacked risk level values.
    """
    columns = data['columns']
    no_risk = data['noRisk']
    low_risk = data['lowRisk']
    medium_risk = data['mediumRisk']
    high_risk = data['highRisk']
    
    # Calculate totals for each month
    totals = [no_risk[i] + low_risk[i] + medium_risk[i] + high_risk[i] for i in range(len(columns))]
    
    return {
        'categories': columns,
        'noRisk': no_risk,
        'lowRisk': low_risk,
        'mediumRisk': medium_risk,
        'highRisk': high_risk,
        'total': totals
    }


def _create_resolution_times_chart_data(severity):
    """Create CategoryChartData for resolution times chart"""
    # Get aggregated data from data model
    aggregated_data = security_dashboard_resolution_times_portfolio_data.chart_resolution_times_by_severity(severity)
    
    # Build chart-ready arrays
    chart_arrays = _build_resolution_times_chart_data_arrays(aggregated_data)
    
    # Get legend labels from data model
    labels = security_dashboard_resolution_times_portfolio_data.get_legend_labels(severity)
    
    # Create CategoryChartData
    chart_data = CategoryChartData()
    chart_data.categories = chart_arrays['categories']
    
    # Add series: All will be stacked, Total displays on top
    chart_data.add_series(labels['noRisk'], chart_arrays['noRisk'])
    chart_data.add_series(labels['lowRisk'], chart_arrays['lowRisk'])
    chart_data.add_series(labels['mediumRisk'], chart_arrays['mediumRisk'])
    chart_data.add_series(labels['highRisk'], chart_arrays['highRisk'])
    chart_data.add_series('Total', chart_arrays['total'])
    
    return chart_data


class SecurityDashboardCriticalResolutionTimesChartPlaceholder(Placeholder):
    """PowerPoint chart depicting resolution times of critical security findings over the last 12 months.
    
    Stacked bar chart showing distribution across different risk categories based on resolution time.
    """
    
    key = "PORTFOLIO_SECURITY_FINDINGS_CRITICAL"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _create_resolution_times_chart_data("CRITICAL")

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        _populate_chart(presentation, key, value_cb, "PORTFOLIO_SECURITY_RESOLUTION_CRITICAL")


class SecurityDashboardHighResolutionTimesChartPlaceholder(Placeholder):
    """PowerPoint chart depicting resolution times of high security findings over the last 12 months.
    
    Stacked bar chart showing distribution across different risk categories based on resolution time.
    """
    
    key = "PORTFOLIO_SECURITY_FINDINGS_HIGHMED"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _create_resolution_times_chart_data("HIGH")

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        _populate_chart(presentation, key, value_cb, "PORTFOLIO_SECURITY_RESOLUTION_HIGH")


class SecurityDashboardMediumResolutionTimesChartPlaceholder(Placeholder):
    """PowerPoint chart depicting resolution times of medium security findings over the last 12 months.
    
    Stacked bar chart showing distribution across different risk categories based on resolution time.
    """
    
    key = "PORTFOLIO_SECURITY_FINDINGS_HIGHMED"
    __doc_type__ = PlaceholderDocType.CHART

    @classmethod
    def value(cls, parameter=None):
        return _create_resolution_times_chart_data("MEDIUM")

    @staticmethod
    def resolve_pptx(presentation: Presentation, key: str, value_cb: Callable) -> None:
        _populate_chart(presentation, key, value_cb, "PORTFOLIO_SECURITY_RESOLUTION_MEDIUM")

