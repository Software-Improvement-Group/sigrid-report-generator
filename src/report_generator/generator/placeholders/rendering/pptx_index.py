#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

"""Builds and caches the searchable index of a presentation's text locations.

Separated from the rendering API in `pptx` because it answers a different question: `pptx` knows
how to read and write PowerPoint content, this module knows where that content is. Locating a
placeholder used to walk the whole presentation, once per placeholder key; this module walks it
once and answers every later lookup from the result.
"""

from dataclasses import dataclass, field

from pptx.enum.shapes import MSO_SHAPE_TYPE

# noinspection PyProtectedMember
from pptx.text.text import _Paragraph

from . import traversal_cache


@dataclass
class ParagraphRecord:
    """One indexable text location.

    `text_owner` is the proxy whose text is searched, which is not always the paragraph: a table
    cell is matched on the whole cell text but yields its first paragraph. `top_level_shape` is
    recorded during the walk because it cannot be recovered afterwards -- the parent chain of a
    nested paragraph leads to the innermost shape, not to the slide-level one.
    """

    paragraph: _Paragraph
    text_owner: object
    text: str
    top_level_shape: object


@dataclass
class SlideIndex:
    slide: object
    paragraphs: list[ParagraphRecord]
    shape_names_including_nested: set[str]
    top_level_shapes: list


@dataclass
class PresentationIndex:
    """Every text location in a presentation, plus the maps lookups need to stay cheap.

    `paragraphs` is the slide indexes concatenated, so it is in document order. The maps are
    keyed on lxml elements rather than proxies because proxies are recreated on every access
    and so cannot be compared by identity.
    """

    slides: list[SlideIndex]
    paragraphs: list[ParagraphRecord]
    slide_index_by_element: dict[object, SlideIndex]
    records_by_paragraph_element: dict[object, list[ParagraphRecord]]


def _is_graphic_frame(shape) -> bool:
    """A table or chart frame, which has no text frame of its own."""
    return "GraphicFrame" in type(shape).__name__


@dataclass
class _ShapeWalk:
    """The text locations and shape names found under one or more slide-level shapes."""

    records: list[ParagraphRecord] = field(default_factory=list)
    shape_names: set[str] = field(default_factory=set)
    top_level_shape: object = None

    @classmethod
    def under(cls, shape) -> "_ShapeWalk":
        walk = cls()
        walk.descend_from(shape)
        return walk

    def descend_from(self, top_level_shape) -> None:
        self.top_level_shape = top_level_shape
        self.descend(top_level_shape)

    def descend(self, shape) -> None:
        """Deliberately descends through the proxy API rather than the underlying XML: several
        placeholders locate their target by walking back up `paragraph._parent`, and those
        parent links only exist on proxies that were built by descending the same way.
        """
        self.shape_names.add(shape.name.strip())
        if _is_graphic_frame(shape):
            self._collect_table_cells(shape)
        else:
            self._collect_paragraphs(shape)
            self._descend_into_group(shape)

    def _collect_table_cells(self, shape) -> None:
        if not shape.has_table:
            return
        for cell in shape.table.iter_cells():
            self._record(cell.text_frame.paragraphs[0], text_owner=cell)

    def _collect_paragraphs(self, shape) -> None:
        if not shape.has_text_frame:
            return
        for paragraph in shape.text_frame.paragraphs:
            self._record(paragraph, text_owner=paragraph)

    def _descend_into_group(self, shape) -> None:
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            return
        for nested_shape in shape.shapes:
            self.descend(nested_shape)

    def _record(self, paragraph, text_owner) -> None:
        self.records.append(
            ParagraphRecord(
                paragraph=paragraph,
                text_owner=text_owner,
                text=text_owner.text,
                top_level_shape=self.top_level_shape,
            )
        )


def _slide_index(slide) -> SlideIndex:
    top_level_shapes = list(slide.shapes)
    walk = _ShapeWalk()
    for shape in top_level_shapes:
        walk.descend_from(shape)
    return SlideIndex(
        slide=slide,
        paragraphs=walk.records,
        shape_names_including_nested=walk.shape_names,
        top_level_shapes=top_level_shapes,
    )


def _records_by_paragraph_element(
    records: list[ParagraphRecord],
) -> dict[object, list[ParagraphRecord]]:
    grouped: dict[object, list[ParagraphRecord]] = {}
    for record in records:
        # noinspection PyProtectedMember
        grouped.setdefault(record.paragraph._p, []).append(record)
    return grouped


def _presentation_index(presentation) -> PresentationIndex:
    slides = [_slide_index(slide) for slide in presentation.slides]
    paragraphs = [record for slide in slides for record in slide.paragraphs]
    return PresentationIndex(
        slides=slides,
        paragraphs=paragraphs,
        slide_index_by_element={slide.slide.element: slide for slide in slides},
        records_by_paragraph_element=_records_by_paragraph_element(paragraphs),
    )


def for_presentation(presentation) -> PresentationIndex:
    return traversal_cache.index_for(
        presentation, lambda: _presentation_index(presentation)
    )


def for_slide(slide) -> SlideIndex:
    """The cached index for a slide, walking it directly if it is not in the cached document.

    A slide that has been detached from the slide id list is no longer reachable from the
    presentation, so it will not be in the index. Walking it is the honest answer -- returning
    nothing would silently claim the slide has no text.
    """
    presentation = slide.part.package.presentation_part.presentation
    index = traversal_cache.index_for(slide, lambda: _presentation_index(presentation))
    return index.slide_index_by_element.get(slide.element) or _slide_index(slide)


def note_text_changed(paragraph: _Paragraph) -> None:
    """Refresh the cached text of a paragraph that was just written to.

    Cheaper than invalidating: the document structure has not changed, only the text of one
    paragraph, and a later lookup may search for the text that was just written.
    """
    index = traversal_cache.cached_index(paragraph)
    if index is None:
        return
    # noinspection PyProtectedMember
    for record in index.records_by_paragraph_element.get(paragraph._p, ()):
        record.text = record.text_owner.text


def matches(text: str, search_text: str) -> bool:
    """Whether the search text occurs in the text as a whole word, treated literally."""
    return traversal_cache.matches_word_bounded(text, search_text)


def matching_paragraphs(records, search_text):
    return [
        record.paragraph
        for record in records
        if traversal_cache.matches_word_bounded(record.text, search_text)
    ]


def records_including_nested(shape):
    return _ShapeWalk.under(shape).records


def own_records(shape):
    return [
        ParagraphRecord(
            paragraph=paragraph,
            text_owner=paragraph,
            text=paragraph.text,
            top_level_shape=shape,
        )
        for paragraph in shape.text_frame.paragraphs
    ]


def invalidate(anchor) -> None:
    """Drop the cached index because the document structure changed.

    Every helper that removes a slide, a shape or a table row, or adds a paragraph, must call
    this: cached records would otherwise point at detached elements, and writing through them
    fails silently rather than loudly.
    """
    traversal_cache.invalidate(anchor)
