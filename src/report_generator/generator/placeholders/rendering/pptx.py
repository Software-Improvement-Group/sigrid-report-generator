#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import logging
import re
from collections.abc import Iterable
from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.presentation import Presentation

# noinspection PyProtectedMember
from pptx.table import Table, _Row

# noinspection PyProtectedMember
from pptx.text.text import _Paragraph, _Run
from pptx.util import Inches

from report_generator.generator.utils.constants.sentiment import Sentiment

from . import pptx_index
from .common import (
    FontProperties,
    apply_font_properties,
    get_font_properties,
    merge_runs_with_same_formatting,
)


@dataclass
class Hyperlink:
    text: str
    url: str


NA_STAR_COLOR = RGBColor(0x91, 0x90, 0x92)
ONE_STAR_COLOR = RGBColor(0xE0, 0x6C, 0x4F)
TWO_STAR_COLOR = RGBColor(0xE8, 0x99, 0x36)
THREE_STAR_COLOR = RGBColor(0xE9, 0xC3, 0x43)
FOUR_STAR_COLOR = RGBColor(0x68, 0xC0, 0x6B)
FIVE_STAR_COLOR = RGBColor(0x3C, 0x88, 0x42)

SIG_BLUE_COLOR = RGBColor(0x24, 0x35, 0x49)
SIG_GREY_COLOR = RGBColor(0xDF, 0xE2, 0xE7)

RATING_POS_CHANGE_RANGE_COLORS = [RGBColor(0xD9, 0xEE, 0xDD), FIVE_STAR_COLOR]
RATING_NEG_CHANGE_RANGE_COLORS = [RGBColor(0xF3, 0xDD, 0xD7), ONE_STAR_COLOR]

VOLUME_POS_CHANGE_RANGE_COLORS = [
    RGBColor(0xEB, 0xF3, 0xF5),
    RGBColor(0x71, 0xB6, 0xC9),
]
VOLUME_NEG_CHANGE_RANGE_COLORS = [
    RGBColor(0xFA, 0xF1, 0xE1),
    RGBColor(0xE8, 0x99, 0x36),
]

DASHBOARD_EXISTING_FINDINGS_COLOR = RGBColor(0xB5, 0xC4, 0xFF)
DASHBOARD_NEW_FINDINGS_COLOR = RGBColor(0x2E, 0x6B, 0xFF)
DASHBOARD_RESOLVED_FINDINGS_COLOR = RGBColor(0x40, 0xC3, 0x60)

DASHBOARD_RESOLUTION_NO_RISK_COLOR = DASHBOARD_RESOLVED_FINDINGS_COLOR
DASHBOARD_RESOLUTION_LOW_RISK_COLOR = RGBColor(0x3A, 0xA4, 0x98)
DASHBOARD_RESOLUTION_MEDIUM_RISK_COLOR = RGBColor(0x34, 0x8A, 0xC7)
DASHBOARD_RESOLUTION_HIGH_RISK_COLOR = DASHBOARD_NEW_FINDINGS_COLOR


def print_slide_ids(slide):
    # Print slide IDs and names for debugging purposes
    logging.debug("Placeholders:")
    for shape in slide.placeholders:
        logging.debug(f"{shape.placeholder_format.idx} {shape.name}")
    logging.debug("----\n")
    logging.debug("Shapes:")
    for shape in slide.shapes:
        logging.debug(
            f"{shape.shape_id} [{shape.name}] {'(This is a chart)' if shape.has_chart else ''}"
        )


def update_many_paragraphs(
    paragraphs, placeholder_id, replacement_text, font: FontProperties = None
):
    for paragraph in paragraphs:
        update_paragraph(paragraph, placeholder_id, replacement_text, font)


def update_paragraph(
    paragraph: _Paragraph, placeholder_id, replacement_text, font: FontProperties = None
):
    merge_runs_with_same_formatting(paragraph)

    try:
        run_with_placeholder = next(
            run
            for run in (paragraph.runs or [])
            if re.search(rf"\b{re.escape(placeholder_id)}\b", run.text)
        )
    except StopIteration:
        logging.warning(
            f"Attempt to update placeholder '{placeholder_id}', but not found in paragraph: {paragraph.text}"
        )
        return

    logging.debug(
        f'Replacing: {placeholder_id} with "{replacement_text}". New text: {run_with_placeholder.text}'
    )
    run_with_placeholder.text = re.sub(
        rf"\b{re.escape(placeholder_id)}\b",
        str(replacement_text),
        run_with_placeholder.text,
    )

    if font:
        apply_font_properties(run_with_placeholder, font)

    pptx_index.note_text_changed(paragraph)


def _shapes_for_paragraphs(paragraphs):
    # A paragraph is typically in a TextGroup which is in a Shape, so we call getparent() twice.
    # Multiple matching paragraphs can share a shape, so de-duplicate while preserving order.
    shapes = []
    for paragraph in paragraphs:
        # noinspection PyProtectedMember
        shape = paragraph._parent._parent
        if shape not in shapes:
            shapes.append(shape)
    return shapes


def find_shapes_with_text(presentation, search_text):
    shapes = []
    for slide in presentation.slides:
        shapes += _shapes_for_paragraphs(find_text_in_slide(slide, search_text))
    logging.debug(f"Finds for {search_text}: {len(shapes)} shapes")
    return shapes


def find_shapes_with_text_in_slide(slide, search_text):
    return _shapes_for_paragraphs(find_text_in_slide(slide, search_text))


def find_text_in_presentation(presentation, search_text):
    paragraphs = pptx_index.matching_paragraphs(
        pptx_index.for_presentation(presentation).paragraphs, search_text
    )
    logging.debug(f"Finds for {search_text}: {len(paragraphs)} paragraphs")
    return paragraphs


def find_text_in_slide(slide, search_text):
    return pptx_index.matching_paragraphs(
        pptx_index.for_slide(slide).paragraphs, search_text
    )


def find_text_in_table(shape, search_text):
    if not shape.has_table:
        return []
    return pptx_index.matching_paragraphs(
        pptx_index.records_including_nested(shape), search_text
    )


def find_text_in_text_frame(shape, search_text):
    if not shape.has_text_frame:
        return []
    return pptx_index.matching_paragraphs(pptx_index.own_records(shape), search_text)


def find_text_in_shape(shape, search_text):
    return pptx_index.matching_paragraphs(
        pptx_index.records_including_nested(shape), search_text
    )


def add_content_paragraph(text_frame, markers, content, paragraph=None):
    pptx_index.invalidate(text_frame)
    if paragraph is None:
        paragraph = text_frame.add_paragraph()
    for marker in markers:
        set_sig_marker(paragraph, marker)
    run = paragraph.add_run()
    run.text = " " + content


def set_sig_marker(paragraph, marker):
    run = paragraph.add_run()
    run.text = marker
    run.font.name = "SIGMarker"

    # Red, yellow and green colors are taken from the SIG pptx template Signal colors
    if marker == "-":
        run.font.color.rgb = RGBColor(0xCB, 0x55, 0x45)
    if marker == "=":
        run.font.color.rgb = RGBColor(0xF0, 0xC8, 0x5A)
    if marker == "+":
        run.font.color.rgb = RGBColor(0x77, 0xC6, 0x73)


def add_xml_element(parent_xml, tag, **attrs):
    element = OxmlElement(tag)
    element.attrib.update(attrs)
    parent_xml.append(element)
    return element


def remove_shape(shape) -> None:
    """Remove a shape from the slide it is on.

    Invalidates the cached index first: a detached shape's paragraphs would still look writable,
    but the writes are dropped when the file is saved, so a stale index would turn a visible
    failure into a silently wrong report.
    """
    pptx_index.invalidate(shape)
    element = shape.element
    element.getparent().remove(element)


def set_shape_color(shape, rgb_color):
    shape.fill.fore_color.rgb = rgb_color


@dataclass
class ShapeProperties:
    """Visual properties to apply to a pptx shape after text replacement."""

    color: RGBColor
    width_inches: float | None = None
    # After text replacement, shapes may retain the width of the original placeholder key,
    # which is typically longer than the display value. width_inches overrides this to the
    # intended size; width_anchor specifies whether the left or right edge remains fixed.
    width_anchor_right: bool = False


def apply_shape_properties(shape, props: ShapeProperties) -> None:
    set_shape_color(shape, props.color)
    if props.width_inches is not None:
        new_width = Inches(props.width_inches)
        if props.width_anchor_right:
            shape.left += shape.width - new_width
        shape.width = new_width


def determine_rating_color(rating):
    if rating < 0.1:
        return NA_STAR_COLOR
    if rating < 1.5:
        return ONE_STAR_COLOR
    elif rating < 2.5:
        return TWO_STAR_COLOR
    elif rating < 3.5:
        return THREE_STAR_COLOR
    elif rating < 4.5:
        return FOUR_STAR_COLOR
    else:
        return FIVE_STAR_COLOR


SENTIMENT_COLORS = {
    Sentiment.NEGATIVE: ONE_STAR_COLOR,  # red
    Sentiment.NEUTRAL: SIG_BLUE_COLOR,  # blue
    Sentiment.POSITIVE: FIVE_STAR_COLOR,  # green
}


def sentiment_color(sentiment: Sentiment) -> RGBColor:
    return SENTIMENT_COLORS[sentiment]


def test_code_ratio_color(ratio):
    if ratio <= 0.01:
        return ONE_STAR_COLOR
    elif ratio <= 0.15:
        return TWO_STAR_COLOR
    elif ratio <= 0.5:
        return THREE_STAR_COLOR
    elif ratio <= 1.5:
        return FOUR_STAR_COLOR
    else:
        return FIVE_STAR_COLOR


def _slide_title(slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is None or not title_shape.has_text_frame:
        return "<untitled>"
    return title_shape.text_frame.text.strip() or "<untitled>"


def _slide_contains_key(slide, key: str) -> bool:
    # Text placeholders match on paragraph text; chart/table placeholders match on shape name
    # (charts/tables are located by shape.name, see find_charts / find_tables). Both paths
    # descend into group shapes so a placeholder nested in a group is still detected.
    if find_text_in_slide(slide, key):
        return True
    return key in pptx_index.for_slide(slide).shape_names_including_nested


def delete_slides_with_placeholder(
    presentation: Presentation, key: str, reason: str | None = None
) -> None:
    """Remove every slide that contains the given placeholder key.

    Called when a placeholder fails to resolve: the slide is dropped so the report never shows
    an unresolved template token. Idempotent across repeated failures — a slide already removed
    from the id list no longer appears in presentation.slides. `reason`, when given, is the
    message the Sigrid API returned with the failing request (e.g. a missing-license notice).
    """
    id_lst = presentation.slides.element  # <p:sldIdLst>
    detail = f" ({reason})" if reason else ""
    # presentation.slides iterates in <p:sldIdLst> order, so the i-th slide corresponds to the
    # i-th <p:sldId>. Snapshot both to lists up front so removing from id_lst mid-loop is safe.
    removed_any = False
    for slide, sld_id in zip(list(presentation.slides), list(id_lst), strict=True):
        if _slide_contains_key(slide, key):
            id_lst.remove(sld_id)
            removed_any = True
            logging.info(
                f"Skipped slide '{_slide_title(slide)}' containing placeholder '{key}' "
                f"because it failed to resolve{detail}"
            )

    # A report missing a licence fails hundreds of placeholders without matching a single slide,
    # and rebuilding the index for each of those would undo the caching.
    if removed_any:
        pptx_index.invalidate(presentation)


def find_charts(presentation: Presentation, key: str):
    """Find charts by shape name. This is the recommended way to locate charts in a presentation."""
    charts = [
        shape.chart
        for slide in pptx_index.for_presentation(presentation).slides
        for shape in slide.top_level_shapes
        if shape.has_chart and shape.name.strip() == key
    ]
    logging.debug(f"Finds for {key}: {len(charts)}")
    return charts


def find_tables(presentation: Presentation, key: str):
    tables = [
        shape.table
        for slide in pptx_index.for_presentation(presentation).slides
        for shape in slide.top_level_shapes
        if shape.has_table and shape.name.strip() == key
    ]
    logging.debug(f"Finds for {key}: {len(tables)}")
    return tables


def find_shapes(presentation: Presentation, key: str):
    """Find the slide-level shapes whose text contains the key.

    Yields the top-level shape even when the text sits in a nested one, because callers use it
    for the shape's position and size on the slide.
    """
    index = pptx_index.for_presentation(presentation)
    shapes = _deduplicate_consecutive_shapes(
        record.top_level_shape
        for record in index.paragraphs
        if pptx_index.matches(record.text, key)
    )
    logging.debug(f"Finds for {key}: {len(shapes)}")
    return shapes


def _deduplicate_consecutive_shapes(shapes):
    """Drop repeats of the shape that precedes them, comparing underlying elements.

    Every paragraph of a shape is indexed before the next shape's, so a shape that matches
    several times -- a group matching through more than one of its children, say -- always
    produces adjacent entries. Elements are compared because proxies are recreated per access.
    """
    unique = []
    for shape in shapes:
        if not unique or unique[-1].element is not shape.element:
            unique.append(shape)
    return unique


def remove_row_from_table(table: Table, row: _Row):
    pptx_index.invalidate(table)
    # noinspection PyProtectedMember
    tbl = table._tbl
    # noinspection PyProtectedMember
    tr = row._tr
    tbl.remove(tr)


def remove_rows_from_table(table: Table, row_numbers: Iterable[int]):
    reversed_numbers = sorted(row_numbers, reverse=True)
    for row_number in reversed_numbers:
        row = table.rows[row_number]
        remove_row_from_table(table, row)


def update_table(table: Table, value: list[list[str | int | float | Hyperlink]]):
    """
    Fills a PowerPoint table with provided values. Copies formatting from existing cells and applies it to all later cells in that column.
    """
    column_fonts = {}

    for row_idx, row in enumerate(table.rows):
        if row_idx >= len(value):
            remove_rows_from_table(table, range(row_idx, len(table.rows)))
            continue

        for col_idx, cell in enumerate(row.cells):
            if col_idx >= len(value[row_idx]):
                continue

            paragraph: _Paragraph = cell.text_frame.paragraphs[0]
            if paragraph.runs:
                column_fonts[col_idx] = get_font_properties(paragraph.runs[0])

            replace_paragraph_with_text(
                paragraph, value[row_idx][col_idx], column_fonts.get(col_idx)
            )


def _apply_hyperlink(run: _Run, hyperlink: Hyperlink) -> None:
    run.text = hyperlink.text
    run.hyperlink.address = hyperlink.url


def replace_paragraph_with_text(
    paragraph: _Paragraph,
    text: str | int | float | Hyperlink,
    font: FontProperties = None,
):
    paragraph.clear()
    run: _Run = paragraph.add_run()
    if isinstance(text, Hyperlink):
        _apply_hyperlink(run, text)
    else:
        run.text = "" if text is None else str(text)
    if font:
        apply_font_properties(run, font)

    pptx_index.note_text_changed(paragraph)


def interpolate_color(colors, t):
    # Map t to position in color list
    position = t * (len(colors) - 1)
    index = int(position)  # lower bound index
    frac = position - index  # fraction between colors

    # If exactly at the last color
    if index >= len(colors) - 1:
        return colors[-1]

    # Interpolate between the two colors
    r = int(colors[index][0] + (colors[index + 1][0] - colors[index][0]) * frac)
    g = int(colors[index][1] + (colors[index + 1][1] - colors[index][1]) * frac)
    b = int(colors[index][2] + (colors[index + 1][2] - colors[index][2]) * frac)

    # Convert back to hex
    return RGBColor(r, g, b)
