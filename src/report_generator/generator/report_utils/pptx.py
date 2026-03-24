#  Copyright Software Improvement Group
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

import logging
import re
from typing import Iterable, Union

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.presentation import Presentation
# noinspection PyProtectedMember
from pptx.table import Table, _Row
# noinspection PyProtectedMember
from pptx.text.text import _Paragraph, _Run

from .common import FontProperties, apply_font_properties, get_font_properties, merge_runs_with_same_formatting

NA_STAR_COLOR = RGBColor(0x91, 0x90, 0x92)
ONE_STAR_COLOR = RGBColor(0xE0, 0x6C, 0x4F)
TWO_STAR_COLOR = RGBColor(0xE8, 0x99, 0x36)
THREE_STAR_COLOR = RGBColor(0xE9, 0xC3, 0x43)
FOUR_STAR_COLOR = RGBColor(0x68, 0xC0, 0x6B)
FIVE_STAR_COLOR = RGBColor(0x3C, 0x88, 0x42)

SIG_BLUE_COLOR = RGBColor(0x24, 0x35, 0x49)
SIG_GREY_COLOR = RGBColor(0xDF, 0xE2, 0xE7)

MAINTAINABILITY_POS_CHANGE_RANGE_COLORS = [RGBColor(0xD9, 0xEE, 0xDD), FIVE_STAR_COLOR]
MAINTAINABILITY_NEG_CHANGE_RANGE_COLORS = [RGBColor(0xF3, 0xDD, 0xD7), ONE_STAR_COLOR]

VOLUME_POS_CHANGE_RANGE_COLORS = [RGBColor(0xEB, 0xF3, 0xF5), RGBColor(0x71, 0xB6, 0xC9)]
VOLUME_NEG_CHANGE_RANGE_COLORS = [RGBColor(0xFA, 0xF1, 0xE1), RGBColor(0xE8, 0x99, 0x36)]

DASHBOARD_EXISTING_FINDINGS_COLOR = RGBColor(0xB5, 0xC4, 0xFF)
DASHBOARD_NEW_FINDINGS_COLOR = RGBColor(0x2E, 0x6B, 0xFF)
DASHBOARD_RESOLVED_FINDINGS_COLOR = RGBColor(0x40, 0xC3, 0x60)

DASHBOARD_RESOLUTION_NO_RISK_COLOR = DASHBOARD_RESOLVED_FINDINGS_COLOR
DASHBOARD_RESOLUTION_LOW_RISK_COLOR = RGBColor(0x3A, 0xA4, 0x98)
DASHBOARD_RESOLUTION_MEDIUM_RISK_COLOR = RGBColor(0x34, 0x8A, 0xC7)
DASHBOARD_RESOLUTION_HIGH_RISK_COLOR = DASHBOARD_NEW_FINDINGS_COLOR

def print_slide_ids(slide):
    # Print slide IDs and names for debugging purposes
    logging.debug("Placeholders:")
    for shape in slide.placeholders:
        logging.debug('%d %s' % (shape.placeholder_format.idx, shape.name))
    logging.debug("----\n")
    logging.debug("Shapes:")
    for shape in slide.shapes:
        logging.debug('%d [%s] %s' % (shape.shape_id, shape.name, "(This is a chart)" if shape.has_chart else ""))


def update_many_paragraphs(paragraphs, placeholder_id, replacement_text, font: FontProperties = None):
    for paragraph in paragraphs:
        update_paragraph(paragraph, placeholder_id, replacement_text, font)


def update_paragraph(paragraph: _Paragraph, placeholder_id, replacement_text, font: FontProperties = None):
    merge_runs_with_same_formatting(paragraph)

    try:
        run_with_placeholder = next(run for run in (paragraph.runs or []) if re.search(rf'\b{re.escape(placeholder_id)}\b', run.text))
    except StopIteration:
        logging.warning(
            f"Attempt to update placeholder '{placeholder_id}', but not found in paragraph: {paragraph.text}")
        return

    logging.debug(f"Replacing: {placeholder_id} with \"{replacement_text}\". New text: {run_with_placeholder.text}")
    run_with_placeholder.text = re.sub(rf'\b{re.escape(placeholder_id)}\b', str(replacement_text), run_with_placeholder.text)

    if font:
        apply_font_properties(run_with_placeholder, font)


def find_shapes_with_text(presentation, search_text):
    shapes = []
    for slide in presentation.slides:
        paragraphs = find_text_in_slide(slide, search_text)
        # A paragraph is typically in a TextGroup which is in a Shape, so we call getparent() twice
        # noinspection PyProtectedMember
        shapes += [paragraph._parent._parent for paragraph in paragraphs]
    return shapes

def find_shapes_with_text_in_slide(slide, search_text):
    shapes = []
    paragraphs = find_text_in_slide(slide, search_text)
    # A paragraph is typically in a TextGroup which is in a Shape, so we call getparent() twice
    shapes += [paragraph._parent._parent for paragraph in paragraphs]
    return shapes

def find_text_in_presentation(presentation, search_text):
    paragraphs = []
    for slide in presentation.slides:
        paragraphs.extend(find_text_in_slide(slide, search_text))

    return paragraphs


def find_text_in_slide(slide, search_text):
    paragraphs = []
    for shape in slide.shapes:
        result = find_text_in_shape(shape, search_text)
        if result:
            paragraphs.append(result)
    return paragraphs


def find_text_in_table(shape, search_text):
    if shape.has_table:
        for cell in shape.table.iter_cells():
            if re.match(fr".*\b{search_text}\b.*", cell.text):
                return cell.text_frame.paragraphs[0]
    return None


def find_text_in_text_frame(shape, search_text):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if re.match(fr".*\b{search_text}\b.*", paragraph.text):
                return paragraph
    return None


def find_text_in_group(shape, search_text):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            result = find_text_in_shape(s, search_text)
            if result:
                return result
    return None


def find_text_in_shape(shape, search_text):
    if "GraphicFrame" in type(shape).__name__:
        result = find_text_in_table(shape, search_text)
        if result:
            return result
        return None

    result = find_text_in_text_frame(shape, search_text)
    if result:
        return result

    return find_text_in_group(shape, search_text)


def add_content_paragraph(text_frame, markers, content, paragraph=None):
    if paragraph is None:
        paragraph = text_frame.add_paragraph()
    for marker in markers:
        set_sig_marker(paragraph, marker)
    run = paragraph.add_run()
    run.text = " " + content


def set_sig_marker(paragraph, marker):
    run = paragraph.add_run()
    run.text = marker
    run.font.name = "SIGMarker"

    # Red, yellow and green colors are taken from the SIG pptx template Signal colors
    if marker == "-":
        run.font.color.rgb = RGBColor(0xcb, 0x55, 0x45)
    if marker == "=":
        run.font.color.rgb = RGBColor(0xf0, 0xc8, 0x5a)
    if marker == "+":
        run.font.color.rgb = RGBColor(0x77, 0xc6, 0x73)


def add_xml_element(parent_xml, tag, **attrs):
    element = OxmlElement(tag)
    element.attrib.update(attrs)
    parent_xml.append(element)
    return element


def set_shape_color(shape, rgb_color):
    shape.fill.fore_color.rgb = rgb_color


def identify_specific_slide(presentation, marker):
    specific_slides = []
    for slide in presentation.slides:
        if find_text_in_slide(slide, marker):
            specific_slides.append(slide)
    return specific_slides


def determine_rating_color(rating):
    if rating < 0.1:
        return NA_STAR_COLOR
    if rating < 1.5:
        return ONE_STAR_COLOR
    elif rating < 2.5:
        return TWO_STAR_COLOR
    elif rating < 3.5:
        return THREE_STAR_COLOR
    elif rating < 4.5:
        return FOUR_STAR_COLOR
    else:
        return FIVE_STAR_COLOR


def test_code_ratio_color(ratio):
    if ratio <= 0.01:
        return ONE_STAR_COLOR
    elif ratio <= 0.15:
        return TWO_STAR_COLOR
    elif ratio <= 0.5:
        return THREE_STAR_COLOR
    elif ratio <= 1.5:
        return FOUR_STAR_COLOR
    else:
        return FIVE_STAR_COLOR


def gather_charts(presentation: Presentation, key: str):
    """Deprecated but kept for backward compatibility, use find_charts instead so it's not linked to a text box in the slide."""
    charts = []
    for slide in identify_specific_slide(presentation, key):
        for shape in slide.shapes:
            if shape.has_chart:
                charts.append(shape.chart)
    return charts


def find_charts(presentation: Presentation, key: str):
    """Find charts by shape name. This is the recommended way to locate charts in a presentation."""
    return [
        shape.chart
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_chart and shape.name == key
    ]


def find_tables(presentation: Presentation, key: str):
    return [
        shape.table
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_table and shape.name == key
    ]


def find_shapes(presentation: Presentation, key: str):
    return [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if find_text_in_shape(shape, key)
    ]


def remove_row_from_table(table: Table, row: _Row):
    # noinspection PyProtectedMember
    tbl = table._tbl
    # noinspection PyProtectedMember
    tr = row._tr
    tbl.remove(tr)


def remove_rows_from_table(table: Table, row_numbers: Iterable[int]):
    reversed_numbers = sorted(row_numbers, reverse=True)
    for row_number in reversed_numbers:
        row = table.rows[row_number]
        remove_row_from_table(table, row)


def update_table(table: Table, value: list[list[Union[str, int, float]]]):
    """
    Fills a PowerPoint table with provided values. Copies formatting from existing cells and applies it to all later cells in that column.
    """
    column_fonts = {}

    for row_idx, row in enumerate(table.rows):
        if row_idx >= len(value):
            remove_rows_from_table(table, range(row_idx, len(table.rows)))
            continue

        for col_idx, cell in enumerate(row.cells):
            if col_idx >= len(value[row_idx]):
                continue

            paragraph: _Paragraph = cell.text_frame.paragraphs[0]
            if paragraph.runs:
                column_fonts[col_idx] = get_font_properties(paragraph.runs[0])

            replace_paragraph_with_text(paragraph, value[row_idx][col_idx], column_fonts.get(col_idx))


def replace_paragraph_with_text(paragraph: _Paragraph, text: Union[str, int, float], font: FontProperties = None):
    paragraph.clear()

    run: _Run = paragraph.add_run()
    run.text = "" if text is None else str(text)

    if font:
        apply_font_properties(run, font)

def interpolate_color(colors, t):
    # Map t to position in color list
    position = t * (len(colors) - 1)
    index = int(position)           # lower bound index
    frac = position - index         # fraction between colors
    
    # If exactly at the last color
    if index >= len(colors) - 1:
        return colors[-1]
    
    # Interpolate between the two colors
    r = int(colors[index][0] + (colors[index+1][0] - colors[index][0]) * frac)
    g = int(colors[index][1] + (colors[index+1][1] - colors[index][1]) * frac)
    b = int(colors[index][2] + (colors[index+1][2] - colors[index][2]) * frac)
    
    # Convert back to hex
    return RGBColor(r, g, b)